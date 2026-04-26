from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 색상 팔레트 ────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x3A, 0x6B)   # 제목 배경
BLUE       = RGBColor(0x2E, 0x6D, 0xB4)   # 섹션 헤더
LIGHTBLUE  = RGBColor(0xD6, 0xE8, 0xFF)   # 개발 단계 강조 배경
GREEN      = RGBColor(0x10, 0x7C, 0x10)   # 사용 가이드 포인트
LIGHTGREEN = RGBColor(0xD9, 0xF0, 0xD9)   # 사용 가이드 배경
ORANGE     = RGBColor(0xE6, 0x6C, 0x00)   # 강조/경고
LIGHTORANGE= RGBColor(0xFF, 0xF0, 0xD6)   # 강조 배경
GRAY       = RGBColor(0x60, 0x60, 0x60)   # 본문 서브텍스트
LIGHTGRAY  = RGBColor(0xF2, 0xF2, 0xF2)   # 행 배경
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x1A, 0x1A, 0x1A)
TEAL       = RGBColor(0x00, 0x7A, 0x8A)
LIGHTTEAL  = RGBColor(0xD0, 0xF0, 0xF4)
RED        = RGBColor(0xC0, 0x39, 0x2B)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ── 공통 유틸 ─────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Pt as _Pt
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=BLACK, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

def add_textbox_lines(slide, lines, l, t, w, h,
                       size=14, bold_first=False, color=BLACK,
                       line_spacing=None):
    """lines: list of (text, bold, color, size_override)"""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            txt, bld, clr, sz = item, False, color, size
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
            sz  = item[3] if len(item) > 3 else size
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        if line_spacing:
            from pptx.util import Pt as _Pt
            p.line_spacing = _Pt(line_spacing)
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(sz)
        run.font.bold  = bld
        run.font.color.rgb = clr
    return txBox

def slide_header(slide, title, subtitle=None, bg=NAVY, title_color=WHITE, subtitle_color=WHITE):
    """상단 헤더 바 + 제목 — 배경을 먼저 그리고 텍스트를 위에 올림"""
    # 1. 슬라이드 전체 흰 배경
    add_rect(slide, 0, 0, 13.33, 7.5, fill=WHITE)
    # 2. 헤더 바
    add_rect(slide, 0, 0, 13.33, 1.3, fill=bg)
    # 3. 제목 텍스트 (헤더 위)
    add_text(slide, title, 0.4, 0.12, 12.5, 0.7,
             size=28, bold=True, color=title_color, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.82, 12.5, 0.4,
                 size=14, color=subtitle_color, align=PP_ALIGN.LEFT, italic=True)

def footer(slide, page_num, total=36):
    add_rect(slide, 0, 7.1, 13.33, 0.4, fill=NAVY)
    add_text(slide, "STM32G4 회로도 검증 & 펌웨어 자동 생성 Agent",
             0.3, 7.12, 10, 0.3, size=10, color=WHITE)
    add_text(slide, f"{page_num} / {total}", 12.3, 7.12, 1, 0.3,
             size=10, color=WHITE, align=PP_ALIGN.RIGHT)

def phase_badge(slide, label, l, t, bg=BLUE, fg=WHITE, w=1.8, h=0.35):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, label, l+0.05, t+0.02, w-0.1, h-0.04,
             size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)

def bullet_box(slide, items, l, t, w, h, bg=LIGHTBLUE, title=None, title_bg=BLUE):
    add_rect(slide, l, t, w, h, fill=bg)
    offset = 0
    if title:
        add_rect(slide, l, t, w, 0.38, fill=title_bg)
        add_text(slide, title, l+0.1, t+0.04, w-0.2, 0.32,
                 size=13, bold=True, color=WHITE)
        offset = 0.42
    lines = []
    for item in items:
        if isinstance(item, str):
            lines.append((item, False, BLACK, 13))
        else:
            lines.append(item)
    add_textbox_lines(slide, lines, l+0.15, t+offset+0.05, w-0.25, h-offset-0.1, size=13)

# ══════════════════════════════════════════════════════════════
# 슬라이드 1 — 타이틀
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 2.8, 13.33, 2.2, fill=BLUE)

add_text(s, "STM32G4 회로도 검증 &", 1.0, 1.1, 11.0, 1.0,
         size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "펌웨어 자동 생성 Agent", 1.0, 2.0, 11.0, 0.85,
         size=36, bold=True, color=RGBColor(0xA8, 0xD0, 0xFF), align=PP_ALIGN.CENTER)

add_text(s, "개발 계획 · 모델 선택 · 사용 가이드", 1.0, 3.0, 11.0, 0.6,
         size=22, color=WHITE, align=PP_ALIGN.CENTER, italic=True)

add_text(s, "대상 칩 : STM32G4 계열   |   서버 : NVIDIA DGX Spark 128GB   |   2026.04",
         1.0, 3.7, 11.0, 0.5, size=14, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.CENTER)

add_rect(s, 0, 6.9, 13.33, 0.6, fill=RGBColor(0x0F, 0x24, 0x45))
add_text(s, "사내 전용 · 외부망 차단 환경", 0.3, 6.95, 12.5, 0.4,
         size=12, color=RGBColor(0x99, 0xBB, 0xFF), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 2 — 목차
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "목차")
footer(s, 2)

sections = [
    ("01", "프로젝트 개요",           "무엇을 만드는가, 왜 만드는가",                  3),
    ("02", "전체 아키텍처",           "3-Step 파이프라인 구성",                       4),
    ("03", "시스템 구성",             "서버, 모델, 도구 스택",                        5),
    ("04", "개발 로드맵",             "Phase 1~5 타임라인 개요",                      6),
    ("05", "Phase별 개발 상세",       "데이터 구축 → 검증 → CubeMX → 통합 → 학습",   "7~11"),
    ("06", "학습 데이터 전략",        "BLDC 알고리즘 모듈 & 데이터 수집 방법",        12),
    ("07", "모델 선택 가이드",        "Step별 LLM 선택 기준 & 메모리 배치",           13),
    ("08", "사용 가이드",             "회로도 입력 → 검증 → 코드생성 → 통합 → 빌드", "14~18"),
    ("09", "기대 효과",               "정확도, 시간 절감, 확장 계획",                  19),
]

for i, (num, title, desc, page) in enumerate(sections):
    row = i // 2; col = i % 2
    lx = 0.4 + col * 6.5; ty = 1.5 + row * 1.15
    add_rect(s, lx, ty, 6.1, 1.0, fill=LIGHTBLUE if col==0 else LIGHTTEAL)
    add_text(s, num, lx+0.1, ty+0.08, 0.6, 0.5, size=24, bold=True, color=BLUE)
    add_text(s, title, lx+0.75, ty+0.06, 5.1, 0.45, size=16, bold=True, color=BLACK)
    add_text(s, desc,  lx+0.75, ty+0.5,  5.1, 0.4,  size=12, color=GRAY)
    add_text(s, f"p.{page}", lx+4.9, ty+0.06, 1.1, 0.35, size=11, color=BLUE, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════════
# 슬라이드 3 — 프로젝트 개요
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "프로젝트 개요", "무엇을 자동화하는가")
footer(s, 3)

add_text(s, "문제 정의", 0.4, 1.4, 4.0, 0.4, size=15, bold=True, color=BLUE)
problems = [
    "HW 엔지니어가 회로도 설계 후 FW 개발자에게 핀맵 전달",
    "FW 개발자가 STM32G4 AF 스펙 수동 확인 → 오류 발생 빈번",
    "HAL 초기화 코드 수작업 작성 → 반복 공수 낭비",
    "모터 제어 알고리즘(FOC/전류제어) 통합 코드 매번 재작성",
]
for i, p in enumerate(problems):
    add_rect(s, 0.4, 1.85+i*0.55, 5.9, 0.48, fill=LIGHTORANGE)
    add_text(s, f"▶  {p}", 0.55, 1.87+i*0.55, 5.6, 0.44, size=12.5, color=BLACK)

add_text(s, "해결 방안", 6.8, 1.4, 6.0, 0.4, size=15, bold=True, color=GREEN)
solutions = [
    ("Step 1", "핀 검증 Agent",     "LLM이 AF 스펙 자동 검증 → 오류 리포트"),
    ("Step 2", "CubeMX 자동화",     ".ioc 자동 생성 → CubeMX CLI → HAL 코드"),
    ("Step 3", "알고리즘 통합",      "기존 FOC 코드를 USER CODE 영역에 자동 삽입"),
]
for i, (badge, title, desc) in enumerate(solutions):
    ty = 1.85 + i * 0.75
    add_rect(s, 6.8, ty, 6.1, 0.65, fill=LIGHTGREEN)
    add_rect(s, 6.8, ty, 1.1, 0.65, fill=GREEN)
    add_text(s, badge, 6.82, ty+0.1, 1.05, 0.45, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, 8.0, ty+0.04, 4.7, 0.3, size=14, bold=True, color=GREEN)
    add_text(s, desc,  8.0, ty+0.33, 4.7, 0.28, size=11.5, color=GRAY)

# 핵심 원칙 박스
add_rect(s, 0.4, 4.15, 12.5, 0.7, fill=NAVY)
add_text(s, "💡  핵심 원칙 : LLM은 HAL API를 직접 생성하지 않는다."
         "  CubeMX가 보장하는 정확한 코드 위에 알고리즘 레이어만 조립.",
         0.6, 4.22, 12.1, 0.55, size=14, bold=True, color=WHITE)

# 기대 효과
add_text(s, "기대 효과", 0.4, 5.0, 12.0, 0.38, size=15, bold=True, color=BLUE)
effects = ["핀 오류 사전 차단", "HAL 코드 생성\n90% 단축", "알고리즘 통합\n자동화", "멀티모터\n2~4개 지원", "사내 스타일\n일관성 유지"]
for i, e in enumerate(effects):
    lx = 0.4 + i * 2.5
    add_rect(s, lx, 5.42, 2.35, 0.65, fill=LIGHTBLUE)
    add_text(s, e, lx+0.08, 5.44, 2.2, 0.6, size=12, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 4 — 전체 아키텍처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "전체 아키텍처 — 3-Step 파이프라인")
footer(s, 4)

# Step 박스 3개 (폭 좁혀서 게이트 공간 확보)
steps = [
    ("STEP 1", "핀 검증 Agent", BLUE,
     ["회로도 입력\n(CSV/KiCad/PDF)", "CubeMX XML 핀 DB", "규칙 엔진", "STM32G4 RAG"],
     "Gemma-4-31B\n(~40GB)", "검증 완료\n핀 JSON"),
    ("STEP 2", "CubeMX 자동화", TEAL,
     [".ioc 파일 자동 생성", "CubeMX CLI 실행", "HAL 코드 출력"],
     "LLM 불필요\n(결정론적 처리)", "HAL 초기화 코드\n(main.c, tim.c ...)"),
    ("STEP 3", "알고리즘 통합", GREEN,
     ["Golden Module RAG", "USER CODE 삽입", "모듈 간 연결 코드"],
     "Gemma-4-31B\n(~32GB)", "완성 펌웨어\n(motor_foc.c ...)"),
]

for i, (step, title, color, inputs, model, output) in enumerate(steps):
    lx = 0.3 + i * 4.33
    bg = LIGHTBLUE if color==BLUE else (LIGHTTEAL if color==TEAL else LIGHTGREEN)
    # 헤더
    add_rect(s, lx, 1.35, 3.9, 0.48, fill=color)
    add_text(s, step,  lx+0.08, 1.37, 0.9,  0.42, size=12, bold=True, color=WHITE)
    add_text(s, title, lx+1.05, 1.4,  2.75, 0.38, size=13, bold=True, color=WHITE)
    # 본문 박스
    add_rect(s, lx, 1.85, 3.9, 2.75, fill=bg)
    add_text(s, "입력", lx+0.1, 1.9, 0.9, 0.28, size=10.5, bold=True, color=color)
    for j, inp in enumerate(inputs):
        add_text(s, f"• {inp}", lx+0.18, 2.2+j*0.45, 3.55, 0.42, size=11, color=BLACK)
    # 모델 배지
    add_rect(s, lx, 4.63, 3.9, 0.55, fill=color)
    add_text(s, f"모델 : {model}", lx+0.08, 4.66, 3.75, 0.49, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 출력
    add_rect(s, lx, 5.22, 3.9, 0.62, fill=RGBColor(0xE8,0xE8,0xE8))
    add_text(s, f"출력 : {output}", lx+0.08, 5.28, 3.75, 0.52, size=11, color=BLACK, align=PP_ALIGN.CENTER)

# ── Step1→Gate 화살표 + 게이트 박스 ──────────────────────
# Step1 출력 아래 → 게이트
add_text(s, "↓", 2.05, 5.84, 0.5, 0.38, size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 게이트 마름모 (사각형으로 표현)
add_rect(s, 1.0, 6.22, 2.6, 0.52, fill=RGBColor(0xFF, 0xF0, 0x80))  # 노란 게이트
add_rect(s, 1.0, 6.22, 2.6, 0.52, fill=None, line=ORANGE, line_w=Pt(2))
add_text(s, "🔒 검증 게이트", 1.08, 6.26, 2.44, 0.22, size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, "errors[ ] == 0 ?",  1.08, 6.47, 2.44, 0.22, size=10.5, color=BLACK, align=PP_ALIGN.CENTER)

# PASS 분기 → Step2
add_text(s, "PASS ▶", 3.62, 6.35, 0.9, 0.3, size=10.5, bold=True, color=GREEN)

# FAIL 분기 → 아래 차단 박스
add_text(s, "FAIL ↓", 1.25, 6.75, 0.85, 0.28, size=10.5, bold=True, color=RED)
add_rect(s, 0.3, 7.05, 3.0, 0.35, fill=RGBColor(0xFF, 0xE0, 0xE0))
add_rect(s, 0.3, 7.05, 3.0, 0.35, fill=None, line=RED, line_w=Pt(1.5))
add_text(s, "❌ 진행 차단 — 오류 리포트 반환 → 회로도 수정 후 재실행",
         0.35, 7.07, 2.95, 0.3, size=9.5, color=RED, bold=True)

# ── Step2→Step3 화살표 ──────────────────────────────────
add_text(s, "▶", 4.23, 3.2, 0.42, 0.42, size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# 하단 노트 (게이트 오른쪽 영역으로 이동)
add_rect(s, 3.6, 6.22, 9.4, 0.88, fill=LIGHTORANGE)
add_text(s, "※  Step 2 · Step 3는 Step 1 검증 PASS 시에만 실행됩니다.",
         3.75, 6.25, 9.1, 0.3, size=11.5, bold=True, color=ORANGE)
add_text(s, "Step 2는 LLM 없이 동작.  Step 1 · 3 모델 합산 ~72GB → DGX Spark 128GB 내 동시 상주 가능",
         3.75, 6.57, 9.1, 0.5, size=11, color=BLACK)

# ══════════════════════════════════════════════════════════════
# 슬라이드 5 — 시스템 구성
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "시스템 구성", "서버 · 모델 · 도구 스택")
footer(s, 5)

# 서버 스펙
add_rect(s, 0.35, 1.4, 5.8, 2.6, fill=LIGHTBLUE)
add_rect(s, 0.35, 1.4, 5.8, 0.42, fill=BLUE)
add_text(s, "AI 서버 — NVIDIA DGX Spark", 0.5, 1.44, 5.5, 0.35, size=14, bold=True, color=WHITE)
specs = [
    "칩셋 : GB10 Grace Blackwell Superchip",
    "메모리 : 128GB LPDDR5X (CPU+GPU Unified)",
    "AI 성능 : ~1 PFLOPS (FP4)",
    "OS : Ubuntu (로컬, 외부망 차단)",
    "확장 : 2대 NVLink → 256GB (405B 모델 가능)",
]
for i, spec in enumerate(specs):
    add_text(s, f"• {spec}", 0.5, 1.9+i*0.4, 5.4, 0.38, size=12.5, color=BLACK)

# LLM 모델
add_rect(s, 6.55, 1.4, 6.4, 2.6, fill=LIGHTGREEN)
add_rect(s, 6.55, 1.4, 6.4, 0.42, fill=GREEN)
add_text(s, "LLM 모델 (Ollama 로컬 실행)", 6.7, 1.44, 6.1, 0.35, size=14, bold=True, color=WHITE)
models = [
    ("Step 1 — 검증", "Gemma-4-31B-It (Q4_K_M, ~40GB)"),
    ("Step 3 — 통합", "Gemma-4-31B-It (Q8, ~32GB)"),
    ("임베딩",         "BAAI/bge-m3 (한국어+영어 지원)"),
    ("동시 로드",      "~72GB 합산 → 128GB 내 여유 ✓"),
]
for i, (label, val) in enumerate(models):
    add_text(s, label, 6.7, 1.92+i*0.42, 1.8, 0.38, size=11, bold=True, color=GREEN)
    add_text(s, val,   8.55, 1.92+i*0.42, 4.2, 0.38, size=11.5, color=BLACK)

# 도구 스택 테이블
add_rect(s, 0.35, 4.1, 12.6, 0.42, fill=NAVY)
add_text(s, "도구 스택", 0.5, 4.14, 12.0, 0.35, size=14, bold=True, color=WHITE)

tool_rows = [
    ("RAG 프레임워크", "LlamaIndex",       "문서 파이프라인, RAG 파이프라인 구성"),
    ("벡터 DB",        "Qdrant (Docker)",  "메타데이터 필터링 — 칩/모듈 종류별 필터"),
    ("HAL 코드 생성",  "STM32CubeMX CLI",  "핀 JSON → .ioc → HAL 코드 결정론적 생성"),
    ("API 서버",       "FastAPI",          "사내 서비스 연동 엔드포인트"),
    ("UI",             "Streamlit",        "3-Step 진행 화면, 검증 리포트 표시"),
    ("컴파일 검증",    "arm-none-eabi-gcc","생성 코드 자동 빌드 테스트 CI"),
]
for i, (cat, tool, desc) in enumerate(tool_rows):
    bg = WHITE if i % 2 == 0 else LIGHTGRAY
    add_rect(s, 0.35, 4.55+i*0.37, 12.6, 0.37, fill=bg)
    add_text(s, cat,  0.5,  4.58+i*0.37, 2.1, 0.32, size=11.5, bold=True, color=NAVY)
    add_text(s, tool, 2.7,  4.58+i*0.37, 2.5, 0.32, size=11.5, bold=True, color=BLUE)
    add_text(s, desc, 5.35, 4.58+i*0.37, 7.3, 0.32, size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 6 — 개발 로드맵
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "개발 로드맵", "Phase 1~5 타임라인")
footer(s, 6)

phases = [
    ("Phase 1", "데이터 기반 구축",     "4~6주",  BLUE,  "CubeMX DB 파싱\nGolden Module 등록\n기존 FOC 코드 태깅"),
    ("Phase 2", "Step 1\n핀 검증 Agent", "3~4주", TEAL,  "규칙 엔진 구현\nRAG 구성\nGemma-4-31B 연결"),
    ("Phase 3", "Step 2\nCubeMX 자동화", "2~3주", GREEN, ".ioc 생성기\nCubeMX CLI 연동\n출력 검증"),
    ("Phase 4", "Step 3\n통합 Agent",   "4~6주",  ORANGE,"Golden Module RAG\nGemma-4-31B\nUSER CODE 삽입"),
    ("Phase 5", "Fine-tuning\n(선택)",  "8~12주", RED,   "학습 데이터 300개\nQLoRA 파인튜닝\n성능 평가"),
]

total_w = 12.4
for i, (phase, title, duration, color, tasks) in enumerate(phases):
    lx = 0.45 + i * (total_w / len(phases) + 0.05)
    bw = total_w / len(phases) - 0.05

    # 헤더
    add_rect(s, lx, 1.4, bw, 0.42, fill=color)
    add_text(s, phase, lx+0.05, 1.42, bw-0.1, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # 제목
    add_rect(s, lx, 1.83, bw, 0.72, fill=RGBColor(0xEE,0xEE,0xEE))
    add_text(s, title, lx+0.05, 1.85, bw-0.1, 0.68, size=12.5, bold=True, color=color, align=PP_ALIGN.CENTER)
    # 기간 뱃지
    add_rect(s, lx, 2.57, bw, 0.32, fill=color)
    add_text(s, duration, lx+0.05, 2.59, bw-0.1, 0.28, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    # 태스크
    task_lines = tasks.split("\n")
    add_rect(s, lx, 2.92, bw, len(task_lines)*0.48+0.2, fill=RGBColor(0xF8,0xF8,0xF8))
    for j, task in enumerate(task_lines):
        add_text(s, f"• {task}", lx+0.1, 2.98+j*0.48, bw-0.15, 0.44, size=11.5, color=BLACK)

# 연결 화살표
for i in range(4):
    lx = 0.45 + (i+1) * (total_w / len(phases) + 0.05) - 0.12
    add_text(s, "▶", lx, 1.88, 0.3, 0.4, size=14, color=NAVY, bold=True)

# 하단 주석
add_rect(s, 0.45, 5.7, 12.4, 0.55, fill=LIGHTBLUE)
add_text(s,
    "MVP 목표 : Phase 1~3 완료 시 기본 동작 가능 (핀 검증 + CubeMX 코드 생성)  |  "
    "Phase 4 완료 시 알고리즘 통합까지 풀 파이프라인  |  Phase 5는 데이터 축적 후 진행",
    0.6, 5.75, 12.1, 0.45, size=11.5, color=NAVY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 7 — Phase 1: 데이터 기반 구축
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Phase 1 — 데이터 기반 구축", "기간: 4~6주  |  LLM 없이 진행 가능")
footer(s, 7)
phase_badge(s, "Phase 1", 0.4, 1.35, bg=BLUE)

add_text(s, "① CubeMX XML → 핀 AF JSON DB", 0.4, 1.85, 6.0, 0.38, size=14, bold=True, color=BLUE)
tasks1 = [
    "STM32CubeMX\\db\\mcu\\STM32G474RETx.xml 파싱",
    "모든 핀 → Alternate Function 목록 추출",
    "결과: {\"PA8\": [\"TIM1_CH1\", \"I2C2_SCL\", ...]} 형태 JSON",
    "G4 전체 계열 커버 (G431 / G471 / G474 / G491 등)",
]
for i, t in enumerate(tasks1):
    add_text(s, f"    {'▸' if i>0 else '📌'} {t}", 0.4, 2.28+i*0.42, 6.0, 0.4, size=12.5, color=BLACK)

add_text(s, "② 기존 보유 코드 → Golden Module 등록", 0.4, 4.1, 6.0, 0.38, size=14, bold=True, color=BLUE)
modules = [
    "foc_clarke / foc_park / foc_inv_park",
    "foc_svpwm / foc_current_pi (FMAC 옵션)",
    "foc_speed_pi / foc_position_pi",
    "foc_angle_encoder / hall / smo",
    "foc_current_sense (G4 OPAMP 연동)",
]
for i, m in enumerate(modules):
    add_text(s, f"    ✅ {m}", 0.4, 4.52+i*0.38, 6.0, 0.36, size=12, color=GREEN)

# 우측
add_text(s, "③ 신규 작성 모듈", 6.8, 1.85, 6.0, 0.38, size=14, bold=True, color=ORANGE)
new_mods = [
    ("bldc_6step_hall.c", "6-step 전환 로직 (Hall 인터럽트)"),
    ("dc_motor_pid.c",    "H-bridge PWM + PID 제어"),
    ("fdcan_motor_cmd.c", "FDCAN 커맨드 파싱/송신"),
    ("multi_axis_sync.c", "2축/3축 TIM 동기화"),
]
for i, (fname, desc) in enumerate(new_mods):
    add_rect(s, 6.8, 2.28+i*0.58, 6.0, 0.52, fill=LIGHTORANGE)
    add_text(s, f"🔨 {fname}", 6.95, 2.3+i*0.58, 2.8, 0.28, size=12, bold=True, color=ORANGE)
    add_text(s, desc,          9.9,  2.3+i*0.58, 2.7, 0.28, size=11.5, color=GRAY)

add_text(s, "④ 학습 데이터 목표 수량", 6.8, 4.65, 6.0, 0.38, size=14, bold=True, color=BLUE)
add_rect(s, 6.8, 5.08, 2.85, 0.7, fill=LIGHTBLUE)
add_text(s, "RAG 지식베이스\n50~80개 (품질 중심)", 6.9, 5.13, 2.7, 0.6, size=12.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_rect(s, 9.95, 5.08, 2.85, 0.7, fill=LIGHTGREEN)
add_text(s, "Fine-tuning 데이터셋\n300~500개 (다양성)", 10.05, 5.13, 2.7, 0.6, size=12.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 8 — Phase 2: 핀 검증 Agent
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Phase 2 — Step 1 핀 검증 Agent", "기간: 3~4주")
footer(s, 8)
phase_badge(s, "Phase 2", 0.4, 1.35, bg=TEAL)

add_text(s, "구현 항목", 0.4, 1.85, 5.8, 0.38, size=15, bold=True, color=TEAL)
impl_items = [
    ("규칙 엔진",   "핀 AF 호환성 검사, 핀 중복, 전용핀 오용, 전원핀 누락"),
    ("RAG 구성",   "STM32G4 Datasheet PDF → 청킹 → BGE-M3 임베딩 → Qdrant"),
    ("LLM 연결",   "Gemma-4-31B: 검증 결과 자연어 설명 + JSON 구조화 출력"),
    ("UI 구성",    "Streamlit: 회로도 업로드 → 검증 진행 → 리포트 표시"),
]
for i, (label, desc) in enumerate(impl_items):
    add_rect(s, 0.4, 2.3+i*0.72, 5.8, 0.65, fill=LIGHTTEAL)
    add_text(s, label, 0.55, 2.34+i*0.72, 1.4, 0.3, size=12, bold=True, color=TEAL)
    add_text(s, desc,  2.1,  2.34+i*0.72, 3.9, 0.55, size=12, color=BLACK)

add_text(s, "검증 항목", 6.8, 1.85, 6.1, 0.38, size=15, bold=True, color=TEAL)
checks = [
    ("✅", "AF 호환성",  "PA8 → TIM1_CH1 (AF6) ✓  /  USART1_TX ✗"),
    ("✅", "핀 중복",    "동일 핀에 두 기능 배정 감지"),
    ("✅", "전용 핀",    "OSC_IN/OUT, NRST, BOOT0 오용 감지"),
    ("✅", "전원 핀",    "VDD, VDDA, VSS 연결 누락 확인"),
    ("✅", "클럭 일관성","HSE/HSI/LSE 설정과 핀 일치 여부"),
    ("✅", "G4 특화",    "CORDIC/FMAC 사용 시 전용 핀 확인"),
]
for i, (icon, label, desc) in enumerate(checks):
    add_rect(s, 6.8, 2.3+i*0.52, 6.1, 0.48, fill=WHITE if i%2==0 else LIGHTGRAY)
    add_text(s, icon,  6.9,  2.34+i*0.52, 0.35, 0.4, size=14)
    add_text(s, label, 7.3,  2.34+i*0.52, 1.5,  0.4, size=12.5, bold=True, color=TEAL)
    add_text(s, desc,  9.0,  2.34+i*0.52, 3.7,  0.4, size=11.5, color=GRAY)

add_rect(s, 0.4, 5.52, 12.5, 0.58, fill=LIGHTBLUE)
add_text(s, "출력 : 검증 완료 핀 JSON (chip, validation, pins[], peripherals[], errors[], warnings[]) → Step 2 입력",
         0.6, 5.57, 12.1, 0.48, size=12.5, color=NAVY, bold=True)

# ══════════════════════════════════════════════════════════════
# 슬라이드 9 — Phase 3: CubeMX 자동화
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Phase 3 — Step 2  CubeMX 자동화", "기간: 2~3주  |  LLM 불필요 — 완전 결정론적")
footer(s, 9)
phase_badge(s, "Phase 3", 0.4, 1.35, bg=GREEN)

# 흐름
flow_steps = [
    ("핀 JSON 수신", "Step 1\n출력", BLUE),
    (".ioc 생성기", "Python 스크립트\n자동 변환", GREEN),
    ("CubeMX CLI", "STM32CubeMX -q\nscript.txt", TEAL),
    ("HAL 코드 출력", "main.c / tim.c\nadc.c / fdcan.c", ORANGE),
]
for i, (title, desc, color) in enumerate(flow_steps):
    lx = 0.4 + i * 3.15
    add_rect(s, lx, 1.4, 2.9, 0.45, fill=color)
    add_text(s, title, lx+0.1, 1.43, 2.7, 0.38, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, lx, 1.87, 2.9, 0.6, fill=RGBColor(0xEE,0xEE,0xEE))
    add_text(s, desc, lx+0.1, 1.9, 2.7, 0.55, size=12, color=BLACK, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "▶", lx+2.95, 1.87, 0.3, 0.45, size=16, bold=True, color=NAVY)

add_text(s, ".ioc 파일 구조 (텍스트 포맷, 자동 생성)", 0.4, 2.65, 6.2, 0.38, size=14, bold=True, color=GREEN)
ioc_sample = """Mcu.UserName=STM32G474RET6
RCC.HSEFreq_Value=24000000
RCC.SYSCLKFreq_VALUE=170000000
PA8.Signal=TIM1_CH1
PA8.GPIO_Label=PWM_UH
PA7.Signal=TIM1_CH1N
PA7.GPIO_Label=PWM_UL
PB9.Signal=FDCAN1_TX
IP.TIM1=enabled
IP.ADC1=enabled
IP.FDCAN1=enabled
IP.CORDIC=enabled"""
add_rect(s, 0.4, 3.08, 6.0, 3.2, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, ioc_sample, 0.55, 3.13, 5.8, 3.1, size=10.5, color=RGBColor(0xCE, 0xF5, 0xB0))

add_text(s, "CubeMX 출력 구조", 6.8, 2.65, 5.8, 0.38, size=14, bold=True, color=GREEN)
output_tree = """project/
├── Core/Src/
│   ├── main.c         ← USER CODE BEGIN/END 포함
│   ├── tim.c          ← MX_TIM1_Init (PWM + 데드타임)
│   ├── adc.c          ← MX_ADC1_Init (Injected + TIM 트리거)
│   ├── fdcan.c        ← MX_FDCAN1_Init
│   └── cordic.c       ← MX_CORDIC_Init
├── Core/Inc/
│   └── main.h
└── Makefile"""
add_rect(s, 6.8, 3.08, 6.0, 2.6, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, output_tree, 6.95, 3.13, 5.8, 2.5, size=11, color=RGBColor(0xAD, 0xD8, 0xE6))

add_rect(s, 6.8, 5.73, 6.0, 0.52, fill=LIGHTGREEN)
add_text(s, "💡  USER CODE BEGIN/END 영역이 Step 3에서 LLM이 채울 공간",
         6.95, 5.78, 5.7, 0.42, size=12.5, bold=True, color=GREEN)

# ══════════════════════════════════════════════════════════════
# 슬라이드 10 — Phase 4: 알고리즘 통합 Agent
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Phase 4 — Step 3  알고리즘 통합 Agent", "기간: 4~6주  |  Gemma-4-31B")
footer(s, 10)
phase_badge(s, "Phase 4", 0.4, 1.35, bg=ORANGE)

add_text(s, "LLM이 하는 일 (범위 명확히 제한)", 0.4, 1.85, 6.2, 0.38, size=14, bold=True, color=ORANGE)
do_items = [
    "USER CODE 영역 내 함수 호출 및 변수 선언",
    "ADC 완료 인터럽트 → FOC 루틴 연결",
    "FDCAN 수신 콜백 → 커맨드 파싱 호출",
    "메인 루프 → 속도/위치 제어 루프 호출",
    "Golden Module include + 파라미터 설정",
]
for i, item in enumerate(do_items):
    add_rect(s, 0.4, 2.28+i*0.5, 6.0, 0.45, fill=LIGHTGREEN)
    add_text(s, f"✅  {item}", 0.55, 2.32+i*0.5, 5.8, 0.38, size=12, color=BLACK)

dont_items = [
    "HAL_TIM_PWM_Init 등 HAL 초기화 함수",
    "Clarke/Park/SVPWM 알고리즘 내부",
    "CORDIC/FMAC 레지스터 직접 접근",
]
add_text(s, "LLM이 하지 않는 일", 0.4, 4.85, 6.2, 0.38, size=13, bold=True, color=RED)
for i, item in enumerate(dont_items):
    add_text(s, f"❌  {item}", 0.55, 5.25+i*0.38, 6.0, 0.35, size=12, color=RED)

add_text(s, "Golden Module RAG 구성", 6.8, 1.85, 6.1, 0.38, size=14, bold=True, color=ORANGE)
rag_items = [
    ("검색 대상",   "기존 FOC 알고리즘 11개 모듈\n+ 신규 작성 4개 모듈"),
    ("임베딩 모델", "BAAI/bge-m3 (코드+주석 혼재)"),
    ("벡터 DB",     "Qdrant + 메타데이터 필터\n(motor_type, control_method 필터)"),
    ("검색 방식",   "Hybrid: Dense + BM25 키워드"),
]
for i, (label, val) in enumerate(rag_items):
    add_rect(s, 6.8, 2.28+i*0.72, 6.1, 0.65, fill=LIGHTORANGE)
    add_text(s, label, 6.95, 2.32+i*0.72, 1.6, 0.28, size=11.5, bold=True, color=ORANGE)
    add_text(s, val,   8.65, 2.32+i*0.72, 4.1, 0.55, size=11.5, color=BLACK)

add_rect(s, 6.8, 5.23, 6.1, 1.02, fill=LIGHTBLUE)
add_rect(s, 6.8, 5.23, 6.1, 0.35, fill=BLUE)
add_text(s, "통합 후 데이터 흐름 (FOC 루프)", 6.95, 5.26, 5.8, 0.3, size=12, bold=True, color=WHITE)
add_text(s,
    "ADC ISR → foc_clarke() → foc_park() → foc_current_pi()\n"
    "→ foc_inv_park() → foc_svpwm() → TIM1 CCR 갱신",
    6.95, 5.62, 5.8, 0.6, size=11.5, color=NAVY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 11 — Phase 5: Fine-tuning
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "Phase 5 — Fine-tuning (선택)", "기간: 8~12주  |  DGX Spark 단독 실행 가능")
footer(s, 11)
phase_badge(s, "Phase 5", 0.4, 1.35, bg=RED)

add_text(s, "목적 & 조건", 0.4, 1.85, 5.8, 0.35, size=14, bold=True, color=RED)
conditions = [
    "Phase 1~4 완료 후 RAG 품질이 부족하다고 판단될 때 진행",
    "학습 데이터 300개 이상 확보된 시점",
    "사내 코딩 스타일 / 네이밍 규칙을 LLM에 체화",
]
for i, c in enumerate(conditions):
    add_text(s, f"• {c}", 0.55, 2.28+i*0.45, 5.6, 0.42, size=12.5, color=BLACK)

add_text(s, "모델별 Fine-tuning 계획", 0.4, 3.62, 5.8, 0.35, size=14, bold=True, color=RED)
ft_rows = [
    ("검증 모델", "Gemma-4-31B", "QLoRA", "~60~80GB", "검증 리포트 포맷 통일"),
    ("통합 모델", "Gemma-4-31B", "QLoRA", "~30~40GB", "사내 코딩 스타일 적용"),
]
add_rect(s, 0.4, 4.05, 5.8, 0.35, fill=NAVY)
for j, hdr in enumerate(["모델 역할", "베이스 모델", "방법", "메모리", "목적"]):
    add_text(s, hdr, 0.5+j*1.1, 4.08, 1.05, 0.28, size=10.5, bold=True, color=WHITE)
for i, row in enumerate(ft_rows):
    bg = WHITE if i%2==0 else LIGHTGRAY
    add_rect(s, 0.4, 4.42+i*0.45, 5.8, 0.42, fill=bg)
    for j, val in enumerate(row):
        add_text(s, val, 0.5+j*1.1, 4.46+i*0.45, 1.05, 0.35, size=11, color=BLACK)

# 우측 — 데이터 수집 전략
add_text(s, "학습 데이터 수집 전략", 6.8, 1.85, 6.1, 0.35, size=14, bold=True, color=RED)
strategies = [
    ("1순위", "기존 FOC 코드 + CubeMX 생성 코드 조합",     "즉시 시작 가능, ~20~30개"),
    ("2순위", "ST MCSDK 생성 코드 (다양한 파라미터 조합)", "~30~50개 추가 확보"),
    ("3순위", "신규 Golden Module 조합 자동 생성",          "컴파일 검증 후 확장"),
    ("품질",  "arm-none-eabi-gcc 빌드 통과 필수",           "전체 샘플 100% 적용"),
    ("검증",  "실제 하드웨어 동작 확인",                   "전체의 20~30% 샘플링"),
]
for i, (rank, desc, note) in enumerate(strategies):
    add_rect(s, 6.8, 2.28+i*0.62, 6.1, 0.56, fill=LIGHTORANGE if i<3 else LIGHTGRAY)
    add_text(s, rank, 6.92, 2.32+i*0.62, 0.85, 0.28, size=11, bold=True, color=RED if i<3 else GRAY)
    add_text(s, desc, 7.85, 2.32+i*0.62, 4.9, 0.28, size=11.5, bold=True, color=BLACK)
    add_text(s, note, 7.85, 2.58+i*0.62, 4.9, 0.24, size=10.5, color=GRAY)

add_rect(s, 0.4, 5.65, 12.5, 0.52, fill=LIGHTBLUE)
add_text(s, "💡  DGX Spark 128GB → A100 없이 72B / 32B QLoRA 파인튜닝 모두 가능",
         0.6, 5.7, 12.1, 0.42, size=13, bold=True, color=NAVY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 12 — 학습 데이터 & BLDC 알고리즘 구조
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "학습 데이터 전략 — BLDC 제어 알고리즘 구조", "기존 보유 최적화 코드 → Golden Module 직접 등록")
footer(s, 12)

add_text(s, "FOC 제어 루프 캐스케이드 (STM32G4 가속기 포함)", 0.4, 1.42, 7.0, 0.38, size=14, bold=True, color=BLUE)

loop_steps = [
    ("위치 PI\n(외부 루프, ~1kHz)", LIGHTGRAY, GRAY),
    ("속도 PI\n(중간 루프, ~1kHz)", LIGHTBLUE, BLUE),
    ("Id/Iq PI\n전류 제어\n(20kHz)", LIGHTGREEN, GREEN),
    ("SVPWM\n→ TIM1 PWM", LIGHTORANGE, ORANGE),
]
for i, (label, bg, fc) in enumerate(loop_steps):
    lx = 0.4 + i * 1.65
    add_rect(s, lx, 1.88, 1.5, 1.15, fill=bg)
    add_text(s, label, lx+0.05, 1.95, 1.4, 1.0, size=11, bold=True, color=fc, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, "→", lx+1.54, 2.3, 0.18, 0.4, size=14, bold=True, color=NAVY)

# 감지 소스
add_rect(s, 0.4, 3.12, 6.5, 0.38, fill=NAVY)
add_text(s, "각도 소스 (θ)", 0.55, 3.15, 3.0, 0.3, size=12, bold=True, color=WHITE)
sensors = [("증분 엔코더", BLUE), ("Hall 센서", TEAL), ("Sensorless SMO", GREEN)]
for i, (name, color) in enumerate(sensors):
    add_rect(s, 0.4+i*2.17, 3.52, 2.1, 0.42, fill=color)
    add_text(s, name, 0.45+i*2.17, 3.55, 2.05, 0.35, size=12, color=WHITE, align=PP_ALIGN.CENTER)

# G4 가속기
add_rect(s, 0.4, 4.05, 6.5, 0.35, fill=ORANGE)
add_text(s, "STM32G4 가속기 활용", 0.55, 4.08, 3.5, 0.28, size=12, bold=True, color=WHITE)
accels = [("CORDIC", "Park/역Park sin/cos", ORANGE), ("FMAC", "Id/Iq PI 필터", ORANGE), ("Internal OPAMP", "전류 감지", ORANGE)]
for i, (hw, use, color) in enumerate(accels):
    add_rect(s, 0.4+i*2.17, 4.43, 2.1, 0.62, fill=LIGHTORANGE)
    add_text(s, hw,  0.5+i*2.17, 4.46, 2.0, 0.28, size=12, bold=True, color=ORANGE)
    add_text(s, use, 0.5+i*2.17, 4.72, 2.0, 0.28, size=11, color=GRAY)

# 우측 — Golden Module 목록
add_text(s, "Golden Module 목록", 7.2, 1.42, 5.8, 0.38, size=14, bold=True, color=BLUE)
gm_rows = [
    ("foc_clarke.c",        "✅", "Clarke 변환 (abc→αβ)"),
    ("foc_park.c",          "✅", "Park 변환, CORDIC 최적화"),
    ("foc_inv_park.c",      "✅", "역 Park 변환"),
    ("foc_svpwm.c",         "✅", "Space Vector PWM"),
    ("foc_current_pi.c",    "✅", "Id/Iq PI + Anti-windup"),
    ("foc_speed_pi.c",      "✅", "속도 PI 제어기"),
    ("foc_position_pi.c",   "✅", "위치 PI 제어기"),
    ("foc_angle_encoder.c", "✅", "증분 엔코더 각도/속도"),
    ("foc_angle_hall.c",    "✅", "Hall 센서 → 각도 매핑"),
    ("foc_angle_smo.c",     "✅", "Sensorless SMO 옵저버"),
    ("foc_current_sense.c", "✅", "ADC + G4 OPAMP 전류감지"),
    ("bldc_6step_hall.c",   "🔨", "6-step 전환 (신규)"),
    ("dc_motor_pid.c",      "🔨", "DC 모터 H-bridge PID (신규)"),
    ("fdcan_motor_cmd.c",   "🔨", "FDCAN 핸들러 (신규)"),
    ("multi_axis_sync.c",   "🔨", "2~4축 TIM 동기화 (신규, 멀티모터)"),
]
for i, (fname, status, desc) in enumerate(gm_rows):
    bg = LIGHTGREEN if status=="✅" else LIGHTORANGE
    add_rect(s, 7.2, 1.85+i*0.33, 5.8, 0.30, fill=bg)
    add_text(s, status, 7.3,  1.87+i*0.33, 0.3,  0.26, size=10)
    add_text(s, fname,  7.65, 1.87+i*0.33, 2.4,  0.26, size=9.5, bold=True, color=GREEN if status=="✅" else ORANGE)
    add_text(s, desc,   10.15,1.87+i*0.33, 2.7,  0.26, size=9.5, color=GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 13 — 모델 선택 가이드
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "모델 선택 가이드", "DGX Spark 128GB 기준 최적 배치")
footer(s, 13)

# Step별 모델 요약
add_rect(s, 0.4, 1.4, 12.5, 0.42, fill=NAVY)
for j, hdr in enumerate(["", "역할", "선택 모델", "양자화", "메모리", "선택 이유"]):
    add_text(s, hdr, 0.45+j*2.08, 1.43, 2.0, 0.35, size=12, bold=True, color=WHITE)

step_rows = [
    ("Step 1", "핀 검증\n(추론)", "Gemma-4-31B-It", "Q4_K_M", "~40GB",
     "한국어+영어 최상급, JSON 구조화 출력 우수"),
    ("Step 2", "CubeMX\n자동화", "LLM 불필요", "—", "0GB",
     "결정론적 코드 생성 — CubeMX CLI 직접 사용"),
    ("Step 3", "알고리즘\n통합", "Gemma-4-31B-It", "Q8", "~32GB",
     "C 코드 생성 최적화, 한국어 주석, USER CODE 삽입"),
    ("임베딩", "RAG 검색", "BAAI/bge-m3", "FP32", "~2GB",
     "한국어+영어 혼재 문서 최적, 로컬 실행 가능"),
]
for i, (step, role, model, quant, mem, reason) in enumerate(step_rows):
    bg = LIGHTBLUE if i%2==0 else WHITE
    add_rect(s, 0.4, 1.85+i*0.68, 12.5, 0.65, fill=bg)
    add_text(s, step,   0.5,  1.9+i*0.68, 0.9, 0.55, size=12.5, bold=True, color=NAVY)
    add_text(s, role,   1.4,  1.9+i*0.68, 1.5, 0.55, size=11.5, color=GRAY)
    add_text(s, model,  2.95, 1.9+i*0.68, 2.8, 0.55, size=11.5, bold=True, color=BLUE)
    add_text(s, quant,  5.8,  1.9+i*0.68, 0.9, 0.55, size=11.5, color=BLACK, align=PP_ALIGN.CENTER)
    add_text(s, mem,    6.75, 1.9+i*0.68, 0.9, 0.55, size=12, bold=True,
             color=GREEN if i!=1 else GRAY, align=PP_ALIGN.CENTER)
    add_text(s, reason, 7.75, 1.9+i*0.68, 4.9, 0.55, size=11, color=GRAY)

# 동시 로드 표
add_text(s, "동시 메모리 로드 시나리오", 0.4, 4.68, 6.2, 0.38, size=14, bold=True, color=BLUE)
load_rows = [
    ("권장",       "Gemma-4-31B Q4", "Gemma-4-31B Q8", "~72GB",  "✓ 여유"),
    ("최고 품질",  "Gemma-4-31B Q8", "Gemma-4-31B Q8", "~104GB", "✓ 가능"),
    ("빠른 응답",  "Gemma-4-31B Q4", "DeepSeek-Coder-V2-Lite", "~56GB","✓ 여유"),
]
add_rect(s, 0.4, 5.1, 9.0, 0.35, fill=NAVY)
for j, h in enumerate(["시나리오", "Step 1 모델", "Step 3 모델", "합산", "판정"]):
    add_text(s, h, 0.5+j*1.8, 5.13, 1.7, 0.28, size=11, bold=True, color=WHITE)
for i, row in enumerate(load_rows):
    bg = LIGHTBLUE if i==0 else (WHITE if i%2==1 else LIGHTGRAY)
    add_rect(s, 0.4, 5.48+i*0.42, 9.0, 0.4, fill=bg)
    for j, val in enumerate(row):
        clr = GREEN if j==4 else (BLUE if j==0 else BLACK)
        add_text(s, val, 0.5+j*1.8, 5.52+i*0.42, 1.7, 0.35, size=11, color=clr, bold=(j==0))

# 우측 — Fine-tuning 옵션
add_text(s, "Fine-tuning 메모리 요구량", 9.8, 4.68, 3.7, 0.38, size=13, bold=True, color=RED)
ft_items = [
    ("Gemma-4-31B QLoRA",       "~60~80GB"),
    ("Gemma-4-31B QLoRA", "~30~40GB"),
]
for i, (m, mem) in enumerate(ft_items):
    add_rect(s, 9.8, 5.1+i*0.58, 3.5, 0.52, fill=LIGHTORANGE)
    add_text(s, m,   9.95, 5.13+i*0.58, 2.4, 0.44, size=11, bold=True, color=ORANGE)
    add_text(s, mem, 12.45,5.13+i*0.58, 0.8, 0.44, size=13, bold=True, color=RED, align=PP_ALIGN.RIGHT)

add_rect(s, 9.8, 6.28, 3.5, 0.42, fill=LIGHTGREEN)
add_text(s, "✅ DGX Spark 단독 가능", 9.95, 6.32, 3.3, 0.34, size=12, bold=True, color=GREEN)

# ══════════════════════════════════════════════════════════════
# 슬라이드 14 — 사용 가이드 개요
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=WHITE)
add_rect(s, 0, 0, 13.33, 1.3, fill=GREEN)
add_text(s, "사용 가이드", 0.4, 0.1, 12.5, 0.75, size=32, bold=True, color=WHITE)
add_text(s, "회로도 입력부터 완성 펌웨어까지 — Step별 사용법", 0.4, 0.85, 12.5, 0.4,
         size=16, color=WHITE, italic=True)
footer(s, 14)

# 전체 흐름
add_text(s, "전체 사용 흐름", 0.4, 1.45, 12.5, 0.38, size=15, bold=True, color=GREEN)

guide_steps = [
    ("①", "준비",           GREEN,  "회로도 핀맵\n작성 (CSV)"),
    ("②", "핀 검증",        GREEN,  "Step 1 Agent\n실행"),
    ("③", "검증 게이트",    ORANGE, "PASS만\nStep 2 진행"),
    ("④", "HAL 코드 생성",  TEAL,   "Step 2\nCubeMX 자동화"),
    ("⑤", "알고리즘 통합",  GREEN,  "Step 3 Agent\n실행"),
    ("⑥", "빌드 확인",      GREEN,  "컴파일 &\n검증"),
]
for i, (num, title, color, desc) in enumerate(guide_steps):
    lx = 0.4 + i * 2.1
    add_rect(s, lx, 1.92, 1.9, 0.52, fill=color)
    add_text(s, f"{num} {title}", lx+0.08, 1.95, 1.75, 0.45, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bg = LIGHTORANGE if color == ORANGE else LIGHTGREEN
    add_rect(s, lx, 2.46, 1.9, 0.65, fill=bg)
    add_text(s, desc, lx+0.08, 2.5, 1.75, 0.57, size=11.5, color=BLACK, align=PP_ALIGN.CENTER)
    if i < 5:
        arrow_color = RED if i == 1 else GREEN  # Step1→게이트 구간 강조
        add_text(s, "▶", lx+1.94, 2.55, 0.22, 0.4, size=14, bold=True, color=arrow_color)

# FAIL 분기 표시 (게이트 아래)
add_rect(s, 2.5, 3.2, 1.9, 0.42, fill=RGBColor(0xFF, 0xE0, 0xE0))
add_rect(s, 2.5, 3.2, 1.9, 0.42, fill=None, line=RED, line_w=Pt(1.5))
add_text(s, "FAIL ↓  오류 리포트", 2.55, 3.23, 1.8, 0.35, size=10.5, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_rect(s, 2.5, 3.65, 1.9, 0.58, fill=RGBColor(0xFF, 0xF0, 0xF0))
add_text(s, "회로도 수정\n후 ① 재실행", 2.55, 3.68, 1.8, 0.52, size=10.5, color=RED, align=PP_ALIGN.CENTER)
# 게이트 박스에서 아래로 화살표
add_text(s, "↓", 3.35, 3.1, 0.35, 0.3, size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)

# 입력 포맷 안내
add_text(s, "입력 파일 포맷 (CSV 예시)", 0.4, 3.3, 6.0, 0.38, size=14, bold=True, color=GREEN)
csv_sample = """chip,STM32G474RET6
hse_mhz,24
sysclk_mhz,170

pin,function,label
PA8,TIM1_CH1,PWM_UH
PA7,TIM1_CH1N,PWM_UL
PA9,TIM1_CH2,PWM_VH
PB0,TIM1_CH2N,PWM_VL
PA10,TIM1_CH3,PWM_WH
PB1,TIM1_CH3N,PWM_WL
PA0,TIM2_CH1,ENC_A
PA1,TIM2_CH2,ENC_B
PB8,FDCAN1_RX,CAN_RX
PB9,FDCAN1_TX,CAN_TX"""
add_rect(s, 0.4, 3.72, 5.8, 3.35, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, csv_sample, 0.55, 3.77, 5.6, 3.25, size=11, color=RGBColor(0xCE,0xF5,0xB0))

# 자연어 프롬프트 입력 가이드
add_text(s, "기능 요구사항 — 자연어 프롬프트 입력", 6.8, 3.3, 6.1, 0.38, size=14, bold=True, color=GREEN)
add_text(s, "HW 개발자가 자유롭게 작성 → LLM이 파싱하여 핀 검증 수행",
         6.8, 3.66, 6.1, 0.24, size=10, color=GRAY, italic=True)
prompt_example = """STM32G474RET6 칩을 쓸 거고,
외부 크리스탈 24MHz / 시스템 170MHz야.

BLDC 모터 1개를 FOC로 제어할 건데
증분형 엔코더(A/B/Z)로 각도 읽고,
3상 6채널 PWM으로 인버터 구동해.
데드타임은 500ns, 전류는 내부 OPAMP.

통신은 FDCAN 1Mbps 쓰고,
파라미터 저장용으로 SPI EEPROM도
연결할 거야. (CS: PA4)"""
add_rect(s, 6.8, 3.93, 6.1, 2.7, fill=RGBColor(0x1A, 0x2A, 0x1A))
add_rect(s, 6.8, 3.93, 0.07, 2.7, fill=GREEN)
add_text(s, prompt_example, 6.97, 3.98, 5.85, 2.6, size=11, color=RGBColor(0xCE, 0xF5, 0xB0))
# 포함 항목 힌트
add_rect(s, 6.8, 6.67, 6.1, 0.4, fill=LIGHTGREEN)
add_text(s, "포함 권장: 칩명 · 클럭 · 모터종류/제어방식 · 피드백센서 · PWM채널수 · 통신프로토콜 · 외부장치(EEPROM 등)",
         6.85, 6.7, 6.0, 0.34, size=9.5, color=RGBColor(0x1A, 0x40, 0x1A))

# ══════════════════════════════════════════════════════════════
# 슬라이드 15 — 사용 가이드: Step 1
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "사용 가이드 ① — Step 1 : 핀 검증", subtitle="회로도 핀맵 CSV + 자연어 프롬프트 입력 → LLM 파싱 → 검증 리포트")
footer(s, 15)
phase_badge(s, "사용 가이드", 0.4, 1.35, bg=GREEN, w=2.2)

add_text(s, "실행 순서", 0.4, 1.85, 5.8, 0.35, size=14, bold=True, color=GREEN)
steps_guide = [
    "1.  웹 UI 접속  (http://dgx-spark:3000)",
    "2.  핀맵 CSV 파일 업로드 (회로도에서 추출)",
    "3.  기능 요구사항을 자연어로 입력 (프롬프트 가이드 참고)",
    "4.  [핀 검증 실행] 버튼 클릭",
    "5.  검증 결과 리포트 확인 (PASS / FAIL)",
]
for i, step in enumerate(steps_guide):
    add_rect(s, 0.4, 2.28+i*0.5, 5.8, 0.46, fill=LIGHTGREEN if i%2==0 else WHITE)
    add_text(s, step, 0.55, 2.32+i*0.5, 5.6, 0.4, size=12.5, color=BLACK)

# 프롬프트 작성 가이드 체크리스트
add_text(s, "프롬프트 작성 가이드 (이런 내용을 포함하세요)", 0.4, 4.83, 5.8, 0.3, size=12, bold=True, color=GREEN)
prompt_items = [
    ("칩 정보",    "어떤 STM32G4 칩인지, 외부 크리스탈, 동작 클럭"),
    ("모터/제어",  "모터 종류(BLDC/DC), 제어 방식(FOC/6-step/듀티), 축 수"),
    ("피드백",     "엔코더(증분/절대), Hall 센서, 리졸버 등"),
    ("PWM/인버터", "채널 수, 데드타임, 전류 감지 방법"),
    ("통신",       "FDCAN, UART, SPI, I2C 및 속도/주소 등"),
    ("외부 장치",  "EEPROM, 센서, 디스플레이 등 연결 장치"),
]
for i, (tag, desc) in enumerate(prompt_items):
    py = 5.16 + i * 0.31
    add_rect(s, 0.4, py, 1.3, 0.27, fill=GREEN)
    add_text(s, tag, 0.42, py+0.02, 1.26, 0.23, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.72, py, 4.48, 0.27, fill=LIGHTGREEN if i%2==0 else WHITE)
    add_text(s, desc, 1.82, py+0.03, 4.3, 0.22, size=10, color=BLACK)

# 리포트 예시
add_text(s, "검증 리포트 출력 예시", 6.8, 1.85, 6.1, 0.35, size=14, bold=True, color=GREEN)
report_ok = """✅ 검증 통과 (PASS)
────────────────────────────
PIN    FUNCTION     AF    STATUS
PA8    TIM1_CH1     AF6   OK
PA7    TIM1_CH1N    AF6   OK
PA9    TIM1_CH2     AF6   OK
PB8    FDCAN1_RX    AF9   OK
PB9    FDCAN1_TX    AF9   OK

⚠ 경고 (Warnings)
• PA0 : 풀업 저항 미설정 권장
• TIM1 데드타임 500ns 설정 필요"""
add_rect(s, 6.8, 2.28, 5.8, 2.85, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, report_ok, 6.95, 2.33, 5.6, 2.75, size=11, color=RGBColor(0xCE,0xF5,0xB0))

report_fail = """❌ 검증 실패 (FAIL)
────────────────────────────
ERROR: PA9 — TIM1_CH2 ≠ AF6
  └ PA9의 TIM1_CH2는 AF6이지만
    회로도에 AF7로 설정됨

ERROR: PB10 — I2C2_SCL 핀 중복
  └ PB10이 FDCAN1_TX와 충돌

→ 회로도 수정 후 재검증 필요"""
add_rect(s, 6.8, 5.2, 5.8, 1.82, fill=RGBColor(0x2A,0x0D,0x0D))
add_text(s, report_fail, 6.95, 5.25, 5.6, 1.72, size=11, color=RGBColor(0xFF,0x9A,0x9A))

add_rect(s, 0.4, 5.7, 5.8, 0.62, fill=LIGHTGREEN)
add_rect(s, 0.4, 5.7, 0.07, 0.62, fill=GREEN)
add_text(s, "✅  PASS :  errors[ ] == 0 → Step 2 자동 진행",
         0.55, 5.73, 5.6, 0.28, size=12.5, bold=True, color=GREEN)
add_text(s, "경고(warnings)는 리포트에 표시하되 진행은 허용",
         0.55, 6.0, 5.6, 0.28, size=11.5, color=GRAY)

add_rect(s, 6.8, 5.7, 6.1, 0.62, fill=RGBColor(0xFF, 0xE8, 0xE8))
add_rect(s, 6.8, 5.7, 0.07, 0.62, fill=RED)
add_text(s, "❌  FAIL :  errors[ ] > 0 → Step 2·3 완전 차단",
         6.95, 5.73, 5.85, 0.28, size=12.5, bold=True, color=RED)
add_text(s, "오류 핀 목록 + 원인 리포트 → 회로도 수정 후 Step 1 재실행",
         6.95, 6.0,  5.85, 0.28, size=11.5, color=GRAY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 16 — 사용 가이드: Step 2
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "사용 가이드 ② — Step 2 : CubeMX 코드 자동 생성", subtitle="Step 1 PASS 확인 후 자동 진행 — 사용자 조작 불필요")
footer(s, 16)
phase_badge(s, "사용 가이드", 0.4, 1.35, bg=GREEN, w=2.2)

add_rect(s, 0.4, 1.85, 12.5, 0.52, fill=TEAL)
add_text(s, "🔒 검증 게이트 통과 후에만 실행됩니다.  errors[ ] > 0 이면 이 화면에 진입 불가.",
         0.6, 1.9, 12.1, 0.42, size=13.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s, "내부 처리 과정 (투명하게 표시)", 0.4, 2.5, 12.5, 0.38, size=14, bold=True, color=TEAL)

process_steps = [
    ("1", "핀 JSON 수신",       "Step 1의 검증 완료 JSON (chip, pins[], peripherals[])"),
    ("2", ".ioc 파일 생성",     "Python 스크립트 → 핀 설정, 클럭, 주변장치 자동 작성"),
    ("3", "CubeMX CLI 실행",    "STM32CubeMX -q script.txt (백그라운드 실행, ~10~30초)"),
    ("4", "출력 파일 수집",     "main.c / tim.c / adc.c / fdcan.c / cordic.c 수집"),
    ("5", "USER CODE 영역 확인","Step 3 LLM이 채울 영역 존재 여부 검증"),
]
for i, (num, title, desc) in enumerate(process_steps):
    add_rect(s, 0.4, 2.95+i*0.65, 0.52, 0.52, fill=TEAL)
    add_text(s, num, 0.4, 2.97+i*0.65, 0.52, 0.48, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 0.95, 2.95+i*0.65, 12.0, 0.52, fill=LIGHTTEAL if i%2==0 else WHITE)
    add_text(s, title, 1.1, 2.99+i*0.65, 2.8, 0.44, size=13, bold=True, color=TEAL)
    add_text(s, desc,  4.05, 2.99+i*0.65, 8.7, 0.44, size=12.5, color=BLACK)

add_rect(s, 0.4, 6.27, 12.5, 0.7, fill=LIGHTGREEN)
add_rect(s, 0.4, 6.27, 0.08, 0.7, fill=GREEN)
add_text(s,
    "출력 확인 :  Streamlit 화면에서 생성된 파일 트리 확인 가능\n"
    "main.c의 USER CODE BEGIN 영역을 미리보기로 표시 → Step 3 진행 확인 후 클릭",
    0.6, 6.32, 12.1, 0.6, size=12.5, color=GREEN)

# ══════════════════════════════════════════════════════════════
# 슬라이드 17 — 사용 가이드: Step 3
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "사용 가이드 ③ — Step 3 : 알고리즘 통합", subtitle="CubeMX 생성 코드 + Golden Module → 완성 펌웨어")
footer(s, 17)
phase_badge(s, "사용 가이드", 0.4, 1.35, bg=GREEN, w=2.2)

add_text(s, "사용 방법", 0.4, 1.85, 5.8, 0.35, size=14, bold=True, color=GREEN)
s3_steps = [
    "1.  Step 2 완료 확인 후 [알고리즘 통합 실행] 버튼 클릭",
    "2.  LLM이 요구사항 JSON을 분석하여 필요한 Golden Module 자동 선택",
    "3.  선택된 모듈 목록 확인 (수동 추가/제거 가능)",
    "4.  [코드 생성 확정] 클릭 → Gemma-4-31B 실행",
    "5.  생성 완료 후 코드 미리보기 확인",
    "6.  [다운로드] 또는 프로젝트 폴더로 저장",
]
for i, step in enumerate(s3_steps):
    add_rect(s, 0.4, 2.28+i*0.5, 5.8, 0.46, fill=LIGHTGREEN if i%2==0 else WHITE)
    add_text(s, step, 0.55, 2.32+i*0.5, 5.6, 0.42, size=12, color=BLACK)

# 선택 모듈 예시
add_text(s, "자동 선택된 Golden Module 예시 (BLDC FOC + FDCAN)", 6.8, 1.85, 6.1, 0.35, size=13, bold=True, color=GREEN)
selected = [
    ("✅", "foc_clarke.c",     "Clarke 변환"),
    ("✅", "foc_park.c",       "Park 변환 (CORDIC)"),
    ("✅", "foc_inv_park.c",   "역 Park 변환"),
    ("✅", "foc_svpwm.c",      "SVPWM"),
    ("✅", "foc_current_pi.c", "Id/Iq PI 제어기"),
    ("✅", "foc_speed_pi.c",   "속도 PI"),
    ("✅", "foc_angle_encoder.c", "엔코더 각도"),
    ("✅", "foc_current_sense.c", "전류 감지"),
    ("✅", "fdcan_motor_cmd.c",   "FDCAN 핸들러"),
]
for i, (icon, fname, desc) in enumerate(selected):
    add_rect(s, 6.8, 2.28+i*0.47, 6.1, 0.44, fill=LIGHTGREEN if i%2==0 else WHITE)
    add_text(s, icon,  6.9,  2.32+i*0.47, 0.3, 0.37, size=12)
    add_text(s, fname, 7.25, 2.32+i*0.47, 2.5, 0.37, size=11.5, bold=True, color=GREEN)
    add_text(s, desc,  9.85, 2.32+i*0.47, 2.9, 0.37, size=11, color=GRAY)

add_rect(s, 0.4, 5.4, 12.5, 0.65, fill=LIGHTBLUE)
add_text(s, "💡  LLM이 삽입하는 코드 :  ADC ISR에 FOC 루틴 호출, FDCAN 콜백에 커맨드 파싱 연결, "
         "메인 루프에 속도 제어 루프 추가 — HAL API는 건드리지 않음",
         0.6, 5.45, 12.1, 0.55, size=12.5, color=NAVY)

# ══════════════════════════════════════════════════════════════
# 슬라이드 18 — 사용 가이드: 결과 확인 & 빌드
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "사용 가이드 ④ — 결과 확인 & 빌드", subtitle="출력 파일 구조 · 빌드 방법 · 주의사항")
footer(s, 18)
phase_badge(s, "사용 가이드", 0.4, 1.35, bg=GREEN, w=2.2)

add_text(s, "최종 출력 파일 구조", 0.4, 1.85, 5.8, 0.35, size=14, bold=True, color=GREEN)
final_tree = """project_output/
├── Core/Src/
│   ├── main.c           ← USER CODE 채워진 통합본
│   ├── tim.c / adc.c / fdcan.c / cordic.c
│   ├── motor_foc.c      ← FOC 알고리즘 통합
│   ├── motor_foc.h
│   └── fdcan_handler.c  ← FDCAN 핸들러
├── Core/Inc/
│   ├── main.h
│   └── motor_params.h   ← 모터 파라미터 설정
├── Drivers/             ← CubeMX HAL 라이브러리
└── Makefile"""
add_rect(s, 0.4, 2.28, 5.8, 3.45, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, final_tree, 0.55, 2.33, 5.6, 3.35, size=11, color=RGBColor(0xAD,0xD8,0xE6))

add_text(s, "빌드 & 검증", 6.8, 1.85, 6.1, 0.35, size=14, bold=True, color=GREEN)
build_steps = [
    ("컴파일",   "make all  (arm-none-eabi-gcc)"),
    ("자동 검증", "CI 파이프라인 → 빌드 결과 Streamlit 표시"),
    ("플래싱",   "make flash  (OpenOCD / ST-Link)"),
    ("디버그",   "make debug  (GDB + OpenOCD)"),
]
for i, (label, cmd) in enumerate(build_steps):
    add_rect(s, 6.8, 2.28+i*0.58, 1.5, 0.52, fill=TEAL)
    add_text(s, label, 6.85, 2.32+i*0.58, 1.4, 0.44, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 8.32, 2.28+i*0.58, 4.58, 0.52, fill=RGBColor(0x1E,0x1E,0x1E))
    add_text(s, cmd, 8.42, 2.32+i*0.58, 4.44, 0.44, size=12, color=RGBColor(0xCE,0xF5,0xB0))

add_text(s, "주의사항 & 권장 작업", 6.8, 4.65, 6.1, 0.35, size=14, bold=True, color=ORANGE)
cautions = [
    "모터 파라미터 (Rs, Ld, Lq, Flux) 실측 후 motor_params.h 수정 필수",
    "FDCAN ID 및 메시지 포맷은 프로젝트 CAN 스펙에 맞게 확인",
    "데드타임 설정값은 사용 드라이버 IC 스펙 확인 후 수정",
    "처음 실행 시 무부하 + 저속으로 FOC 루프 튜닝 권장",
]
for i, c in enumerate(cautions):
    add_rect(s, 6.8, 5.08+i*0.47, 6.1, 0.43, fill=LIGHTORANGE)
    add_text(s, f"⚠  {c}", 6.95, 5.12+i*0.47, 5.9, 0.36, size=11.5, color=BLACK)

# ══════════════════════════════════════════════════════════════
# 슬라이드 19 — GUI 예시 ① 메인 대시보드
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "UI 화면 예시 ① — 메인 대시보드", subtitle="브라우저 접속 초기 화면 — 프로젝트 현황 및 Step 진입")
footer(s, 19)
phase_badge(s, "UI 예시", 0.4, 1.35, bg=TEAL, w=1.8)

# ── UI 컬러 팔레트 ──────────────────────────────────────
_SB   = RGBColor(0x1A, 0x20, 0x2C)   # sidebar bg
_SACT = RGBColor(0x2D, 0x6A, 0xD6)   # sidebar active item
_APB  = RGBColor(0x0F, 0x24, 0x45)   # app top-bar
_PBG  = RGBColor(0xF0, 0xF4, 0xF8)   # page background
_MUT  = RGBColor(0x71, 0x82, 0x96)   # muted text
_THD  = RGBColor(0xED, 0xF2, 0xF7)   # table header row
_ALT  = RGBColor(0xF7, 0xFA, 0xFC)   # alt table row
_PBGC = RGBColor(0xC6, 0xF6, 0xD5)   # PASS badge bg
_PBGF = RGBColor(0x22, 0x54, 0x3D)   # PASS badge fg
_TBGC = RGBColor(0xBE, 0xE3, 0xF8)   # teal badge bg
_TBGF = RGBColor(0x1A, 0x36, 0x5D)   # teal badge fg
_WBGC = RGBColor(0xED, 0xED, 0xED)   # wait badge bg
_WBGF = RGBColor(0x77, 0x77, 0x77)   # wait badge fg

BX, BY = 0.35, 1.85
BW     = 12.63

# 브라우저 크롬 바
add_rect(s, BX, BY, BW, 0.27, fill=RGBColor(0x3C, 0x3C, 0x3C))
add_text(s, "* * *", BX+0.1, BY+0.03, 0.6, 0.21, size=7, color=RGBColor(0xAA,0xAA,0xAA))
add_rect(s, BX+0.85, BY+0.04, 7.6, 0.18, fill=RGBColor(0x50,0x50,0x50))
add_text(s, "http://dgx-spark:3000  —  STM32G4 Agent", BX+0.9, BY+0.04, 7.5, 0.18, size=7.5, color=RGBColor(0xDD,0xDD,0xDD))

# 앱 내비게이션 바
add_rect(s, BX, BY+0.27, BW, 0.38, fill=_APB)
add_text(s, "STM32G4 Agent", BX+0.2, BY+0.32, 3.5, 0.28, size=12, bold=True, color=WHITE)
add_text(s, "대시보드     작업이력     설정     도움말", BX+7.6, BY+0.34, 4.8, 0.24, size=9, color=RGBColor(0x90,0xB4,0xCC))
add_rect(s, BX+11.5, BY+0.3, 0.3, 0.3, fill=TEAL)
add_text(s, "K", BX+11.5, BY+0.3, 0.3, 0.3, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 페이지 본문 영역
PY = BY + 0.65
PH = 7.1 - PY - 0.05
add_rect(s, BX, PY, BW, PH, fill=_PBG)

# 사이드바 (폭 1.9)
SBW = 1.9
add_rect(s, BX, PY, SBW, PH, fill=_SB)
add_text(s, "MENU", BX+0.14, PY+0.14, 1.6, 0.2, size=7.5, bold=True, color=_MUT)
_menus = [
    ("Home  대시보드",     True),
    ("01    핀 검증",      False),
    ("02    코드 생성",    False),
    ("03    알고리즘 통합",False),
    ("Hist  작업 이력",    False),
    ("Cfg   설정",         False),
]
for i, (lbl, act) in enumerate(_menus):
    add_rect(s, BX, PY+0.38+i*0.5, SBW, 0.46, fill=_SACT if act else _SB)
    add_text(s, lbl, BX+0.12, PY+0.43+i*0.5, SBW-0.15, 0.36, size=9.5,
             bold=act, color=WHITE if act else RGBColor(0xA0,0xB4,0xCC))

# 콘텐츠 영역
CX = BX + SBW + 0.05
CW = BW - SBW - 0.1

add_text(s, "대시보드", CX+0.15, PY+0.1, 4.0, 0.3, size=13, bold=True, color=RGBColor(0x1A,0x20,0x2C))
add_text(s, "현재 프로젝트: BLDC_Motor_v2  |  STM32G474RET6", CX+0.15, PY+0.42, 7.5, 0.22, size=9, color=_MUT)

# Step 상태 카드 3개
_CARDS = [
    ("STEP 1", "핀 검증",      "PASS",    "오류 0건  경고 1건",  GREEN, _PBGC, _PBGF),
    ("STEP 2", "HAL 코드 생성","완료",    "파일 5개 생성",       TEAL,  _TBGC, _TBGF),
    ("STEP 3", "알고리즘 통합","대기중",  "Step 2 완료 후 진행", _MUT,  _WBGC, _WBGF),
]
for i, (step, title, status, sub, acc, bb, bf) in enumerate(_CARDS):
    cx2 = CX + 0.15 + i * 3.3
    add_rect(s, cx2, PY+0.68, 3.1, 1.0, fill=WHITE)
    add_rect(s, cx2, PY+0.68, 0.07, 1.0, fill=acc)
    add_text(s, step,   cx2+0.18, PY+0.73, 1.5, 0.22, size=8,  bold=True, color=_MUT)
    add_text(s, title,  cx2+0.18, PY+0.94, 2.8, 0.3,  size=11, bold=True, color=RGBColor(0x1A,0x20,0x2C))
    add_rect(s, cx2+0.18, PY+1.27, 1.05, 0.22, fill=bb)
    add_text(s, status, cx2+0.18, PY+1.28, 1.05, 0.21, size=8.5, bold=True, color=bf, align=PP_ALIGN.CENTER)
    add_text(s, sub,    cx2+0.18, PY+1.53, 2.8,  0.2,  size=8,   color=_MUT)

# 최근 작업 이력 테이블
TY = PY + 1.82
add_text(s, "최근 작업 이력", CX+0.15, TY, 5.0, 0.26, size=11, bold=True, color=RGBColor(0x1A,0x20,0x2C))
TY2 = TY + 0.3
add_rect(s, CX+0.15, TY2, CW-0.2, 0.25, fill=_THD)
_COLS = [("프로젝트명", 2.8), ("칩 모델", 2.0), ("상태", 1.65), ("날짜", 1.7), ("담당자", 1.4)]
rx = CX+0.25
for hdr, cw in _COLS:
    add_text(s, hdr, rx, TY2+0.02, cw-0.05, 0.21, size=8.5, bold=True, color=_MUT)
    rx += cw
_ROWS = [
    ("BLDC_Motor_v2",  "STM32G474RET6", "PASS",     "2026-04-10", "김철수"),
    ("DC_Conv_3kW",    "STM32G431KBT6", "경고 1건", "2026-04-09", "이영희"),
    ("EPS_Controller", "STM32G474VET6", "오류 2건", "2026-04-08", "박민준"),
]
for ri, (p, c, st, dt, u) in enumerate(_ROWS):
    rby = TY2 + 0.25 + ri * 0.33
    add_rect(s, CX+0.15, rby, CW-0.2, 0.31, fill=WHITE if ri%2==0 else _ALT)
    rx2 = CX+0.25
    for v, (_, cw) in zip([p, c, st, dt, u], _COLS):
        clr = GREEN if v=="PASS" else (RED if "오류" in v else (ORANGE if "경고" in v else BLACK))
        add_text(s, v, rx2, rby+0.04, cw-0.05, 0.25, size=9, color=clr)
        rx2 += cw

# 새 프로젝트 버튼
BTN_Y = TY2 + 0.25 + len(_ROWS)*0.33 + 0.2
add_rect(s, CX+0.15, BTN_Y, 2.6, 0.38, fill=BLUE)
add_text(s, "+  새 프로젝트 시작", CX+0.2, BTN_Y+0.05, 2.5, 0.28, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 20 — GUI 예시 ② Step 1 핀 검증 화면
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "UI 화면 예시 ② — Step 1 핀 검증", subtitle="CSV 업로드 + 요구사항 JSON 입력 → 검증 결과 실시간 표시")
footer(s, 20)
phase_badge(s, "UI 예시", 0.4, 1.35, bg=TEAL, w=1.8)

# 브라우저 크롬 + 앱 바 (재사용)
add_rect(s, BX, BY, BW, 0.27, fill=RGBColor(0x3C, 0x3C, 0x3C))
add_text(s, "* * *", BX+0.1, BY+0.03, 0.6, 0.21, size=7, color=RGBColor(0xAA,0xAA,0xAA))
add_rect(s, BX+0.85, BY+0.04, 7.6, 0.18, fill=RGBColor(0x50,0x50,0x50))
add_text(s, "http://dgx-spark:3000/step1", BX+0.9, BY+0.04, 7.5, 0.18, size=7.5, color=RGBColor(0xDD,0xDD,0xDD))
add_rect(s, BX, BY+0.27, BW, 0.38, fill=_APB)
add_text(s, "STM32G4 Agent", BX+0.2, BY+0.32, 3.5, 0.28, size=12, bold=True, color=WHITE)
add_text(s, "대시보드     작업이력     설정", BX+8.3, BY+0.34, 4.0, 0.24, size=9, color=RGBColor(0x90,0xB4,0xCC))
add_rect(s, BX+11.5, BY+0.3, 0.3, 0.3, fill=TEAL)
add_text(s, "K", BX+11.5, BY+0.3, 0.3, 0.3, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 본문
add_rect(s, BX, PY, BW, PH, fill=_PBG)

# 사이드바 (Step 1 active)
add_rect(s, BX, PY, SBW, PH, fill=_SB)
add_text(s, "MENU", BX+0.14, PY+0.14, 1.6, 0.2, size=7.5, bold=True, color=_MUT)
_menus2 = [
    ("Home  대시보드",     False),
    ("01    핀 검증",      True),
    ("02    코드 생성",    False),
    ("03    알고리즘 통합",False),
    ("Hist  작업 이력",    False),
    ("Cfg   설정",         False),
]
for i, (lbl, act) in enumerate(_menus2):
    add_rect(s, BX, PY+0.38+i*0.5, SBW, 0.46, fill=_SACT if act else _SB)
    add_text(s, lbl, BX+0.12, PY+0.43+i*0.5, SBW-0.15, 0.36, size=9.5,
             bold=act, color=WHITE if act else RGBColor(0xA0,0xB4,0xCC))

# 페이지 제목
add_text(s, "STEP 1  |  핀 검증", CX+0.15, PY+0.1, 5.0, 0.3, size=13, bold=True, color=RGBColor(0x1A,0x20,0x2C))

# ── 왼쪽 패널: 입력 폼 ──────────────────────────
LP_X, LP_W = CX+0.1, 4.8

# ① CSV 업로드 존
add_text(s, "1  핀맵 CSV 업로드", LP_X, PY+0.46, LP_W, 0.25, size=10, bold=True, color=NAVY)
add_rect(s, LP_X, PY+0.73, LP_W, 0.82, fill=WHITE)
add_rect(s, LP_X, PY+0.73, LP_W, 0.82, fill=None, line=RGBColor(0x90,0xB4,0xCC), line_w=Pt(1.2))
add_text(s, "[  ]", LP_X+2.1, PY+0.8, 0.6, 0.28, size=14, color=_MUT, align=PP_ALIGN.CENTER)
add_text(s, "CSV 파일을 드래그하거나 클릭하여 업로드", LP_X+0.2, PY+1.07, LP_W-0.4, 0.2, size=8.5, color=_MUT, align=PP_ALIGN.CENTER)
add_text(s, "pinmap_BLDC_v2.csv  (업로드 완료 ✓)", LP_X+0.2, PY+1.27, LP_W-0.4, 0.2, size=8.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
add_rect(s, LP_X, PY+0.73, LP_W, 0.82, fill=None, line=GREEN, line_w=Pt(1.0))

# ② 자연어 프롬프트 입력
add_text(s, "2  기능 요구사항 입력 (자연어)", LP_X, PY+1.64, LP_W, 0.25, size=10, bold=True, color=NAVY)

# 프롬프트 예시 힌트 박스
add_rect(s, LP_X, PY+1.9, LP_W, 0.3, fill=RGBColor(0xE8,0xF4,0xE8))
add_rect(s, LP_X, PY+1.9, 0.07, 0.3, fill=GREEN)
add_text(s, "예시) STM32G474RET6, BLDC FOC, 엔코더, FDCAN 1Mbps, SPI EEPROM...",
         LP_X+0.15, PY+1.93, LP_W-0.2, 0.24, size=8.5, color=RGBColor(0x1A,0x40,0x1A), italic=True)

# 자연어 텍스트 입력 영역
_prompt_text = """STM32G474RET6 칩을 쓸 거고,
외부 크리스탈 24MHz / 시스템 170MHz야.

BLDC 모터 1개를 FOC로 제어할 건데
증분형 엔코더(A/B/Z)로 각도 읽고,
3상 6채널 PWM으로 인버터 구동해.
데드타임 500ns, 전류는 내부 OPAMP.
FDCAN 1Mbps + SPI EEPROM(CS: PA4)."""
add_rect(s, LP_X, PY+2.22, LP_W, 1.3, fill=WHITE)
add_rect(s, LP_X, PY+2.22, LP_W, 1.3, fill=None, line=RGBColor(0x90,0xCF,0x90), line_w=Pt(1.5))
add_text(s, _prompt_text, LP_X+0.1, PY+2.27, LP_W-0.18, 1.22, size=9.5, color=BLACK)

# 검증 실행 버튼
add_rect(s, LP_X, PY+3.6, LP_W, 0.42, fill=GREEN)
add_text(s, "검증 실행  (Step 1 시작)", LP_X+0.1, PY+3.65, LP_W-0.2, 0.32, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── 오른쪽 패널: 검증 결과 ──────────────────────
RP_X = LP_X + LP_W + 0.3
RP_W = CW - LP_W - 0.5

add_text(s, "검증 결과", RP_X, PY+0.46, RP_W, 0.25, size=10, bold=True, color=NAVY)

# PASS 배너
add_rect(s, RP_X, PY+0.73, RP_W, 0.38, fill=_PBGC)
add_rect(s, RP_X, PY+0.73, 0.07, 0.38, fill=GREEN)
add_text(s, "PASS — errors[ ] = 0     경고 1건", RP_X+0.15, PY+0.78, RP_W-0.2, 0.28, size=11, bold=True, color=_PBGF)

# 핀 검증 테이블
add_rect(s, RP_X, PY+1.15, RP_W, 0.25, fill=_THD)
_VCOLS = [("핀", 0.65), ("기능", 1.5), ("AF", 0.65), ("상태", 0.9)]
rvx = RP_X + 0.08
for hdr, cw in _VCOLS:
    add_text(s, hdr, rvx, PY+1.17, cw-0.05, 0.21, size=8, bold=True, color=_MUT)
    rvx += cw
_VROWS = [
    ("PA8",  "TIM1_CH1",    "AF6", "OK"),
    ("PA7",  "TIM1_CH1N",   "AF6", "OK"),
    ("PA9",  "TIM1_CH2",    "AF6", "OK"),
    ("PB8",  "FDCAN1_RX",   "AF9", "OK"),
    ("PB9",  "FDCAN1_TX",   "AF9", "OK"),
    ("PA0",  "TIM2_CH1",    "AF1", "WARNING"),
]
for ri, (pin, fn, af, st) in enumerate(_VROWS):
    rvy = PY + 1.4 + ri * 0.33
    rbg = RGBColor(0xFF,0xFD,0xF0) if st=="WARNING" else (WHITE if ri%2==0 else _ALT)
    add_rect(s, RP_X, rvy, RP_W, 0.31, fill=rbg)
    clr_st = ORANGE if st=="WARNING" else GREEN
    rvx2 = RP_X + 0.08
    for v, (_, cw) in zip([pin, fn, af, st], _VCOLS):
        c = clr_st if v == st else BLACK
        add_text(s, v, rvx2, rvy+0.04, cw-0.05, 0.24, size=9, color=c)
        rvx2 += cw

# 경고 메모
add_rect(s, RP_X, PY+3.4, RP_W, 0.52, fill=RGBColor(0xFF,0xF3,0xCD))
add_rect(s, RP_X, PY+3.4, 0.07, 0.52, fill=ORANGE)
add_text(s, "WARNING — PA0 : 인코더 입력 풀업 저항 미설정 권장\n→ 회로도 수정 권장 (진행은 허용)", RP_X+0.15, PY+3.44, RP_W-0.2, 0.44, size=9, color=RGBColor(0x7A,0x4A,0x00))

# Step 2 진행 버튼
add_rect(s, RP_X, PY+4.02, RP_W, 0.42, fill=TEAL)
add_text(s, "Step 2 로 진행  (HAL 코드 생성)", RP_X+0.1, PY+4.07, RP_W-0.2, 0.32, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 21 — GUI 예시 ③ Step 2/3 화면
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "UI 화면 예시 ③ — Step 2 / Step 3", subtitle="HAL 코드 생성 진행 상황 → 알고리즘 통합 모듈 선택 + 코드 미리보기")
footer(s, 21)
phase_badge(s, "UI 예시", 0.4, 1.35, bg=TEAL, w=1.8)

# ── 상단: Step 2 진행 화면 (좌측 절반) ─────────────
# 작은 레이블
add_text(s, "Step 2 : HAL 코드 자동 생성 (CubeMX CLI)", 0.4, 1.87, 6.2, 0.3, size=11, bold=True, color=TEAL)

# Step 2 패널
S2X, S2Y, S2W, S2H = 0.4, 2.2, 5.9, 2.3
add_rect(s, S2X, S2Y, S2W, S2H, fill=WHITE)
add_rect(s, S2X, S2Y, S2W, 0.35, fill=TEAL)
add_text(s, "코드 생성 완료 — 파일 5개", S2X+0.15, S2Y+0.06, S2W-0.3, 0.26, size=10.5, bold=True, color=WHITE)

_S2STEPS = [
    ("완료", "핀 JSON 수신",       "Step 1 검증 완료 데이터 수신",              GREEN),
    ("완료", ".ioc 파일 생성",     "170MHz 클럭 / TIM1 SVPWM / FDCAN 설정",     GREEN),
    ("완료", "CubeMX CLI 실행",    "STM32CubeMX -q script.txt  (22초)",          GREEN),
    ("완료", "출력 파일 수집",     "main.c  tim.c  adc.c  fdcan.c  cordic.c",    GREEN),
    ("완료", "USER CODE 영역 확인","Step 3 LLM 삽입 영역 5개 확인",              GREEN),
]
for i, (badge, title, desc, bc) in enumerate(_S2STEPS):
    rsy = S2Y + 0.42 + i * 0.36
    add_rect(s, S2X+0.1, rsy, 0.6, 0.3, fill=_PBGC)
    add_text(s, badge, S2X+0.1, rsy+0.04, 0.6, 0.22, size=8, bold=True, color=_PBGF, align=PP_ALIGN.CENTER)
    add_text(s, title, S2X+0.78, rsy+0.04, 1.6, 0.22, size=9.5, bold=True, color=NAVY)
    add_text(s, desc,  S2X+2.45, rsy+0.04, 3.3, 0.22, size=9, color=BLACK)

# 생성된 파일 트리 (우측 절반 상단)
add_text(s, "Step 2 생성 파일 트리 + Step 3 알고리즘 선택", 6.55, 1.87, 6.7, 0.3, size=11, bold=True, color=TEAL)
add_rect(s, 6.55, S2Y, 6.4, S2H, fill=RGBColor(0x1E,0x1E,0x1E))
_tree = """project_output/
├── Core/Src/
│   ├── main.c        ← USER CODE BEGIN 영역 5개
│   ├── tim.c         ← TIM1/2 초기화 완료
│   ├── adc.c         ← ADC1/2 완료
│   ├── fdcan.c       ← FDCAN1 완료
│   └── cordic.c      ← CORDIC 완료
├── Core/Inc/
│   └── main.h
└── Drivers/ (HAL)"""
add_text(s, _tree, 6.7, S2Y+0.08, 6.1, S2H-0.15, size=9.5, color=RGBColor(0xAD,0xD8,0xE6))

# 구분선
add_rect(s, 0.35, 4.6, 12.63, 0.04, fill=RGBColor(0xCB,0xD5,0xE0))
add_text(s, "Step 3 : 알고리즘 통합 (Golden Module 선택 + 코드 삽입)", 0.4, 4.68, 12.0, 0.28, size=11, bold=True, color=GREEN)

# ── Step 3 화면 ─────────────────────────────────────
S3Y = 5.0

# 모듈 선택 패널 (좌)
add_rect(s, 0.4, S3Y, 5.5, 1.95, fill=WHITE)
add_rect(s, 0.4, S3Y, 5.5, 0.32, fill=GREEN)
add_text(s, "Golden Module 자동 선택  (수동 편집 가능)", 0.55, S3Y+0.05, 5.2, 0.24, size=10, bold=True, color=WHITE)
_MODS = [
    ("foc_clarke.c",      True,  "Clarke 변환"),
    ("foc_park.c",        True,  "Park 변환 (CORDIC)"),
    ("foc_inv_park.c",    True,  "역 Park 변환"),
    ("foc_svpwm.c",       True,  "SVPWM"),
    ("foc_current_pi.c",  True,  "Id/Iq PI 제어"),
    ("fdcan_motor_cmd.c", True,  "FDCAN 핸들러"),
]
for i, (fname, checked, desc) in enumerate(_MODS):
    my = S3Y + 0.38 + i * 0.26
    add_rect(s, 0.5, my, 0.22, 0.22, fill=GREEN if checked else WHITE)
    add_rect(s, 0.5, my, 0.22, 0.22, fill=None, line=GREEN, line_w=Pt(1.0))
    add_text(s, "v" if checked else " ", 0.5, my, 0.22, 0.22, size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, fname, 0.8, my+0.02, 2.3, 0.2, size=9, bold=True, color=NAVY)
    add_text(s, desc,  3.2, my+0.02, 2.5, 0.2, size=9, color=_MUT)

# 생성 확정 버튼
add_rect(s, 0.4, S3Y+1.6, 5.5, 0.32, fill=GREEN)
add_text(s, "코드 생성 확정  (Gemma-4-31B 실행)", 0.5, S3Y+1.64, 5.3, 0.24, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# 코드 미리보기 (우)
add_rect(s, 6.1, S3Y, 6.85, 1.95, fill=RGBColor(0x1E,0x1E,0x1E))
add_rect(s, 6.1, S3Y, 6.85, 0.3,  fill=RGBColor(0x33,0x33,0x33))
add_text(s, "main.c — 미리보기", 6.25, S3Y+0.05, 4.0, 0.22, size=9.5, bold=True, color=RGBColor(0xCC,0xCC,0xCC))
_code_prev = """/* USER CODE BEGIN 3 */
  /* FOC main loop - inserted by Agent */
  FOC_Clarke(&i_ab, adc_ia, adc_ib);
  FOC_Park(&i_dq, &i_ab, theta_e);
  FOC_CurrentPI(&v_dq, &i_dq, &ref_dq);
  FOC_InvPark(&v_ab, &v_dq, theta_e);
  FOC_SVPWM(&tim_ccr, &v_ab, vbus);
  __HAL_TIM_SET_COMPARE(&htim1,
    TIM_CHANNEL_1, tim_ccr.u);
/* USER CODE END 3 */"""
add_text(s, _code_prev, 6.2, S3Y+0.35, 6.6, 1.55, size=9, color=RGBColor(0xCE,0xF5,0xB0))

# 다운로드 버튼
add_rect(s, 6.1, S3Y+1.97, 3.3, 0.32, fill=BLUE)
add_text(s, "프로젝트 다운로드 (.zip)", 6.2, S3Y+2.01, 3.1, 0.24, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(s, 9.5, S3Y+1.97, 3.45, 0.32, fill=RGBColor(0x2D,0x3A,0x4A))
add_text(s, "프로젝트 폴더 저장 (서버)", 9.6, S3Y+2.01, 3.25, 0.24, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# 슬라이드 22 — 기대 효과 & 다음 단계
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "기대 효과 & 다음 단계")
footer(s, 22)

effects_data = [
    ("핀 오류 사전 차단",      "규칙 엔진 + LLM 이중 검증\n→ 회로도 수정 공수 절감", BLUE),
    ("HAL 코드 생성 자동화",   "CubeMX CLI 자동화\n→ 초기화 코드 작성 시간 ≈ 0", TEAL),
    ("알고리즘 통합 자동화",   "검증된 Golden Module 조립\n→ FOC 통합 반나절 → 수분", GREEN),
    ("멀티모터 지원 (2~4개)", "TIM1+TIM8+TIM20+HRTIM 자동 배분\n→ FOC×2 / DC×4 구성 자동 생성", ORANGE),
]
add_text(s, "기대 효과", 0.4, 1.4, 12.0, 0.38, size=15, bold=True, color=NAVY)
for i, (title, desc, color) in enumerate(effects_data):
    lx = 0.4 + i * 3.1
    add_rect(s, lx, 1.85, 2.9, 1.35, fill=color)
    add_text(s, title, lx+0.1, 1.92, 2.7, 0.55, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, lx, 2.52, 2.9, 0.68, fill=RGBColor(0xEE,0xEE,0xEE))
    add_text(s, desc, lx+0.1, 2.56, 2.7, 0.6, size=11.5, color=BLACK, align=PP_ALIGN.CENTER)

# 다음 단계
add_text(s, "다음 단계 (즉시 시작 가능한 항목)", 0.4, 3.4, 12.0, 0.38, size=15, bold=True, color=GREEN)
next_steps = [
    ("즉시", "ST 공식 PDF 14종 다운로드 (download_st_docs.sh 실행)",  GREEN),
    ("즉시", "STM32G474 CubeMX XML 파싱 → pin_af_db.json 생성",       GREEN),
    ("1주",  "Qdrant + BGE-M3 RAG 파이프라인 구성 (Phase 1)",         BLUE),
    ("2~3주","dc_motor_pid / multi_axis_sync / bldc_6step 모듈 작성", BLUE),
    ("4~6주","FastAPI Step1 검증 에이전트 + Streamlit MVP",           TEAL),
    ("8주~", "멀티모터 리뷰 로직 + Fine-tuning 데이터 수집",          ORANGE),
]
for i, (when, task, color) in enumerate(next_steps):
    add_rect(s, 0.4+((i%2)*6.2), 3.85+(i//2)*0.62, 0.85, 0.55, fill=color)
    add_text(s, when, 0.42+((i%2)*6.2), 3.89+(i//2)*0.62, 0.82, 0.47,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 1.28+((i%2)*6.2), 3.85+(i//2)*0.62, 5.3, 0.55, fill=LIGHTGRAY)
    add_text(s, task, 1.42+((i%2)*6.2), 3.9+(i//2)*0.62, 5.1, 0.45, size=12, color=BLACK)

# 미결 사항
add_rect(s, 0.4, 6.0, 12.5, 0.62, fill=LIGHTORANGE)
add_text(s,
    "미결 사항 :  ① 회로도 입력 포맷 (CSV + PDF 우선)   "
    "② 멀티모터 OPAMP 부족 시 외부 OPAMP IC 지정 방식   "
    "③ MCSDK 라이선스 확인   ④ Golden Module 담당 엔지니어 지정",
    0.6, 6.05, 12.1, 0.52, size=12, color=ORANGE)


# ══════════════════════════════════════════════════════════════
# 어펜딕스 섹션 표지 — A
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 2.5, 13.33, 2.5, fill=BLUE)
add_text(s, "Appendix A", 1.0, 1.3, 11.0, 0.9,
         size=22, bold=True, color=RGBColor(0xA8,0xD0,0xFF), align=PP_ALIGN.CENTER)
add_text(s, "학습 데이터 수집 가이드", 1.0, 2.65, 11.0, 0.95,
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "무엇을 · 어떻게 · 얼마나 모아야 하는가", 1.0, 3.7, 11.0, 0.6,
         size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 0, 6.9, 13.33, 0.6, fill=RGBColor(0x0F,0x24,0x45))
add_text(s, "Step 1 핀 검증 RAG / Step 1 Fine-tuning / Step 3 통합 RAG / Step 3 Fine-tuning",
         0.5, 6.93, 12.3, 0.4, size=11.5, color=RGBColor(0x99,0xBB,0xFF), align=PP_ALIGN.CENTER)
footer(s, 23)

# ══════════════════════════════════════════════════════════════
# A-1: Step 1 RAG 수집 목록
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "A-1  Step 1 핀 검증 — RAG 지식베이스 수집", "목표: ~1,550 청크")
footer(s, 24)

src_data = [
    ("① CubeMX XML 핀 DB", BLUE,
     "경로: STM32CubeMX\\db\\mcu\\STM32G4*.xml\n파일 수: 40~50개 (G4 전체 패키지)",
     "처리: RAG 불필요 — Python 파싱 → pin_af_db.json\n결과: {chip:{pin:[AF목록]}} 구조화 JSON",
     "규칙 엔진 직접 조회", "100% 정확", LIGHTBLUE),
    ("② STM32G4 Datasheet (DS12288)", TEAL,
     "추출 대상: Table 11~15 (AF 맵) + Table 5~8 (핀 설명)\n대상 페이지: 약 50페이지 / 전체 200p 중",
     "청킹: 핀 테이블 → 행(row) 단위 청크\n일반 설명 → 256토큰 / 32 overlap\n툴: pdfplumber",
     "~300 청크", "벡터 임베딩", LIGHTTEAL),
    ("③ Application Notes (7종) + RM0440", GREEN,
     "AN5306 OPAMP전류센싱 / AN5789 Bootstrap\nAN4277 PWM보호/BRK / AN4539 HRTIM\nAN4220 6Step / AN4835 High-side / AN5036 열관리\n+ RM0440 선택 챕터 (Timer/ADC/FDCAN/CORDIC)",
     "청킹: AN = 회로블록 단위\nRM0440 = 레지스터/기능 설명 단위\n512토큰 / 64 overlap",
     "~1,200 청크", "벡터 임베딩", LIGHTGREEN),
    ("④ Errata Sheet + GitHub 레퍼런스", ORANGE,
     "ES0430 Errata (~40p) — 알려진 핀/주변장치 버그\n+ flatmcu KiCad 회로도 (STM32G473 FOC)\n+ STM32CubeG4 공식 예제 (HRTIM/ADC/OPAMP/FDCAN)",
     "Errata: 이슈 단위 (1이슈=1청크)\nGitHub 코드: 함수 단위 (AST 파싱)",
     "~50 + ~200 청크", "벡터 임베딩", LIGHTORANGE),
]

for i, (title, color, desc1, desc2, count, method, bg) in enumerate(src_data):
    row = i // 2; col = i % 2
    lx = 0.35 + col * 6.5; ty = 1.42 + row * 2.6
    add_rect(s, lx, ty, 6.1, 0.38, fill=color)
    add_text(s, title, lx+0.1, ty+0.04, 5.9, 0.3, size=13, bold=True, color=WHITE)
    add_rect(s, lx, ty+0.38, 6.1, 2.12, fill=bg)
    add_text(s, desc1, lx+0.15, ty+0.44, 5.8, 0.88, size=11.5, color=BLACK)
    add_text(s, desc2, lx+0.15, ty+1.36, 5.8, 0.88, size=11.5, color=GRAY)
    add_rect(s, lx, ty+2.26, 3.0, 0.32, fill=color)
    add_text(s, count, lx+0.1, ty+2.29, 2.85, 0.26, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, lx+3.1, ty+2.26, 3.0, 0.32, fill=RGBColor(0xDD,0xDD,0xDD))
    add_text(s, method, lx+3.2, ty+2.29, 2.85, 0.26, size=11.5, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# A-2: Step 1 Fine-tuning 데이터 자동 생성
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "A-2  Step 1 핀 검증 — Fine-tuning 데이터 자동 생성",
             "규칙 엔진으로 정답 자동 생성 가능 → 대부분 자동화")
footer(s, 25)

add_text(s, "데이터 형식 (입력 → 출력 쌍)", 0.35, 1.42, 5.8, 0.35, size=14, bold=True, color=BLUE)
fmt_code = '''// 입력 (instruction + input)
{
  "chip": "STM32G474RET6",
  "pins": [
    {"pin":"PA8", "function":"TIM1_CH1", "af":"AF6"},
    {"pin":"PA9", "function":"TIM1_CH2", "af":"AF7"} ← 오류
  ]
}

// 출력 (정답)
{
  "validation": "FAIL",
  "errors": [{
    "pin": "PA9",
    "assigned_af": "AF7",
    "correct_af": "AF6",
    "message": "TIM1_CH2는 AF6이어야 합니다"
  }],
  "warnings": [],
  "passed_pins": ["PA8"]
}'''
add_rect(s, 0.35, 1.82, 5.8, 4.3, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, fmt_code, 0.5, 1.87, 5.6, 4.2, size=10, color=RGBColor(0xCE,0xF5,0xB0))

# 우측: 케이스 구성 + 수량
add_text(s, "케이스 구성 및 목표 수량", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=BLUE)
cases = [
    ("① 완전 정상 (PASS)",             "30%", "~600개", GREEN),
    ("② AF 번호 오류 1개",              "25%", "~500개", ORANGE),
    ("③ AF 번호 오류 2개 이상",         "15%", "~300개", ORANGE),
    ("④ 핀 중복 배정",                  "10%", "~200개", RED),
    ("⑤ 전용 핀 오용 (NRST/BOOT0)",    "10%", "~200개", RED),
    ("⑥ 전원 핀 누락 (VDD/VSS)",        "10%", "~200개", RED),
]
for i, (case, pct, cnt, color) in enumerate(cases):
    add_rect(s, 6.65, 1.82+i*0.58, 6.3, 0.52, fill=RGBColor(0xF5,0xF5,0xF5) if i%2==0 else WHITE)
    add_rect(s, 6.65, 1.82+i*0.58, 0.06, 0.52, fill=color)
    add_text(s, case, 6.78, 1.86+i*0.58, 3.5, 0.42, size=12, color=BLACK)
    add_text(s, pct,  10.35,1.86+i*0.58, 0.8, 0.42, size=12, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, cnt,  11.2, 1.86+i*0.58, 1.6, 0.42, size=12, bold=True, color=NAVY, align=PP_ALIGN.RIGHT)

add_rect(s, 6.65, 5.33, 6.3, 0.42, fill=NAVY)
add_text(s, "합계 목표 : 2,000쌍  (자동 생성 3~5일 소요)", 6.78, 5.37, 6.1, 0.34,
         size=13, bold=True, color=WHITE)

add_rect(s, 6.65, 5.8, 6.3, 1.28, fill=LIGHTBLUE)
add_text(s, "자동화 파이프라인", 6.78, 5.84, 6.1, 0.28, size=12, bold=True, color=BLUE)
pipe = ("pin_af_db.json  →  generate_dataset.py  →  train.jsonl\n"
        "• 칩 모델 × 핀 개수 × 케이스 종류 조합 자동 생성\n"
        "• 규칙 엔진으로 정답 검증 → 100% 정확한 레이블\n"
        "• 생성 후 arm-none-eabi-gcc 컴파일 검증 불필요 (핀 검증은 코드 아님)")
add_text(s, pipe, 6.78, 6.16, 6.1, 0.88, size=11.5, color=BLACK)

# ══════════════════════════════════════════════════════════════
# A-3: Step 3 RAG 수집 목록
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "A-3  Step 3 알고리즘 통합 — RAG 지식베이스 수집", "목표: ~2,050 청크")
footer(s, 26)

src3 = [
    ("① Golden Module 소스 코드", BLUE,
     ["foc_clarke / foc_park / foc_inv_park",
      "foc_svpwm / foc_current_pi (FMAC 옵션)",
      "foc_speed_pi / foc_position_pi",
      "foc_angle_encoder / hall / smo",
      "foc_current_sense / bldc_6step / fdcan_motor_cmd",
      "dc_motor_pid / multi_axis_sync (멀티모터)",
      "→ .c + .h 쌍, 총 15개 모듈 (30개 파일)"],
     "함수 단위 청킹 (AST 파싱)\n메타데이터: {module, func, g4_accel, motor_count}\n예상: ~150 청크", LIGHTBLUE),
    ("② STM32CubeG4 공식 예제 (수집 완료)", TEAL,
     ["dataset/opensource/STM32CubeG4/ — 42MB",
      "HRTIM: Basic_Multiple_PWM, CBC_Deadtime",
      "TIM: PWMOutput, BreakAndDeadtime, DMA",
      "ADC: GroupsRegularInjected, OffsetComp",
      "OPAMP: TimerControlMux, PGA, Calibration",
      "FDCAN: Com_IT, Classic_Frame / CORDIC: Sin_DMA"],
     "함수 단위 청킹 (HAL_ API만)\n699개 .c / 462개 .h 파일\n예상: ~200 청크", LIGHTTEAL),
    ("③ CubeMX 자동 생성 예제 모음", GREEN,
     ["30가지 핀 설정 조합으로 CubeMX 실행",
      "TIM1단독 / TIM1+TIM8 / +TIM20",
      "ADC 단순 / Injected+DMA / 멀티채널",
      "FDCAN 있음/없음, CORDIC 있음/없음",
      "→ main.c/tim.c/adc.c/fdcan.c/cordic.c 수집"],
     "파일 단위 + 함수 단위 혼합\n예상: ~500 청크", LIGHTGREEN),
    ("④ Reference Manual (Step 1과 공유)", ORANGE,
     ["Step 1 RAG에서 구축한 RM0440 인덱스 재사용",
      "주변장치 챕터 ~1,200 청크",
      "별도 구축 불필요",
      "Qdrant 컬렉션: stm32g4_knowledge 공유"],
     "Step 1 RAG 재사용\n추가 작업 없음", LIGHTORANGE),
]

for i, (title, color, items, note, bg) in enumerate(src3):
    row = i // 2; col = i % 2
    lx = 0.35 + col * 6.5; ty = 1.42 + row * 2.65
    add_rect(s, lx, ty, 6.1, 0.38, fill=color)
    add_text(s, title, lx+0.1, ty+0.04, 5.9, 0.3, size=13, bold=True, color=WHITE)
    add_rect(s, lx, ty+0.38, 6.1, 2.17, fill=bg)
    for j, item in enumerate(items):
        add_text(s, f"• {item}", lx+0.15, ty+0.44+j*0.33, 5.8, 0.31, size=11, color=BLACK)
    add_rect(s, lx, ty+2.56, 6.1, 0.27, fill=color)
    add_text(s, note.replace('\n', '  |  '), lx+0.1, ty+2.59, 5.9, 0.22,
             size=10.5, bold=True, color=WHITE)

# ══════════════════════════════════════════════════════════════
# A-4: Step 3 Fine-tuning 조합 매트릭스
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "A-4  Step 3 알고리즘 통합 — Fine-tuning 데이터 구성",
             "반자동화 파이프라인  |  목표: 500쌍")
footer(s, 27)

add_text(s, "유효 조합 매트릭스 (~50개 → 조합 자동화로 500개 확장)", 0.35, 1.42, 12.6, 0.35,
         size=14, bold=True, color=BLUE)

matrix_rows = [
    ("BLDC", "FOC",    "1", "엔코더",   "FDCAN",      "5", GREEN),
    ("BLDC", "FOC",    "1", "Hall",     "FDCAN",      "5", GREEN),
    ("BLDC", "FOC",    "1", "Sensorless","FDCAN",     "4", GREEN),
    ("BLDC", "FOC",    "2", "엔코더",   "FDCAN",      "5", GREEN),
    ("BLDC", "FOC",    "2", "Hall",     "FDCAN",      "4", GREEN),
    ("BLDC", "FOC",    "3", "엔코더",   "FDCAN",      "5", BLUE),
    ("BLDC", "6-step", "1", "Hall",     "FDCAN",      "4", TEAL),
    ("BLDC", "6-step", "1", "Hall",     "UART",       "3", TEAL),
    ("DC",   "PID",    "1", "엔코더",   "FDCAN",      "4", ORANGE),
    ("DC",   "PID",    "2", "엔코더",   "FDCAN",      "4", ORANGE),
    ("DC",   "PID",    "3", "엔코더",   "FDCAN",      "3", ORANGE),
    ("BLDC", "FOC",    "1", "엔코더",   "FDCAN+UART", "4", BLUE),
]
hdrs = ["모터 종류", "제어 방식", "축 수", "피드백 센서", "통신", "기본 샘플"]
add_rect(s, 0.35, 1.82, 12.6, 0.38, fill=NAVY)
col_w = [1.8, 1.6, 0.8, 2.0, 2.0, 1.5]
cx = 0.45
for k, (h, w) in enumerate(zip(hdrs, col_w)):
    add_text(s, h, cx, 1.85, w-0.05, 0.3, size=11.5, bold=True, color=WHITE)
    cx += w

for i, row in enumerate(matrix_rows):
    bg = RGBColor(0xF5,0xF5,0xF5) if i%2==0 else WHITE
    add_rect(s, 0.35, 2.22+i*0.33, 12.6, 0.32, fill=bg)
    cx = 0.45
    color = row[6]
    for k, (val, w) in enumerate(zip(row[:6], col_w)):
        bold = (k == 0)
        c = color if k == 0 else BLACK
        add_text(s, val, cx, 2.25+i*0.33, w-0.05, 0.27, size=11.5, bold=bold, color=c)
        cx += w

add_rect(s, 0.35, 6.18, 12.6, 0.32, fill=NAVY)
add_text(s, "기본 50쌍  →  조합 자동화 × 10  →  500쌍 목표  |  컴파일 검증(arm-none-eabi-gcc) 100% 적용",
         0.5, 6.21, 12.3, 0.26, size=12, bold=True, color=WHITE)

# 우측 주석 (생성 파이프라인)
add_rect(s, 0.35, 6.54, 12.6, 0.72, fill=LIGHTBLUE)
pipe_text = ("생성 파이프라인:  CubeMX CLI (Layer1)  +  Golden Module (Layer2)  →  엔지니어 검토  →  "
             "컴파일 검증  →  JSON 저장\n"
             "소요: 기본 50쌍 확보(2~3주) → 자동 확장(스크립트 1~2일) → 최종 검토(1주)")
add_text(s, pipe_text, 0.5, 6.58, 12.3, 0.64, size=11.5, color=NAVY)

# ══════════════════════════════════════════════════════════════
# 어펜딕스 섹션 표지 — B
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 2.5, 13.33, 2.5, fill=TEAL)
add_text(s, "Appendix B", 1.0, 1.3, 11.0, 0.9,
         size=22, bold=True, color=RGBColor(0xA8,0xD0,0xFF), align=PP_ALIGN.CENTER)
add_text(s, "모델 학습 방법 상세", 1.0, 2.65, 11.0, 0.95,
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "환경 설정 · RAG 구축 · QLoRA Fine-tuning · 배포 · 평가", 1.0, 3.7, 11.0, 0.6,
         size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 0, 6.9, 13.33, 0.6, fill=RGBColor(0x0F,0x24,0x45))
add_text(s, "DGX Spark 128GB  |  Unsloth + QLoRA  |  Ollama 배포",
         0.5, 6.93, 12.3, 0.4, size=11.5, color=RGBColor(0x99,0xBB,0xFF), align=PP_ALIGN.CENTER)
footer(s, 28)

# ══════════════════════════════════════════════════════════════
# B-1: 환경 설정 + RAG 구축
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "B-1  환경 설정 & RAG 지식베이스 구축")
footer(s, 29)

add_text(s, "필수 패키지 설치", 0.35, 1.42, 5.8, 0.35, size=14, bold=True, color=TEAL)
pkg_code = """conda create -n stm32_train python=3.11
conda activate stm32_train

# PyTorch (GB10 Blackwell)
pip install torch --index-url \
  https://download.pytorch.org/whl/cu128

# 학습 라이브러리
pip install transformers==4.48.0
pip install peft==0.14.0        # QLoRA
pip install trl==0.13.0         # SFT Trainer
pip install bitsandbytes==0.45.0
pip install unsloth             # 2× 속도 가속
pip install accelerate==1.3.0

# RAG / 벡터 DB
pip install llama-index
pip install llama-index-vector-stores-qdrant
pip install llama-index-embeddings-huggingface
pip install qdrant-client"""
add_rect(s, 0.35, 1.82, 5.8, 3.9, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, pkg_code, 0.5, 1.87, 5.6, 3.8, size=10, color=RGBColor(0xCE,0xF5,0xB0))

add_text(s, "Qdrant 실행 + RAG 인덱싱", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=TEAL)
qdrant_code = """# docker-compose.yml
services:
  qdrant:
    image: qdrant/qdrant:latest
    ports: ["6333:6333"]
    volumes:
      - ./qdrant_storage:/qdrant/storage

# 실행
docker compose up -d qdrant

# RAG 인덱스 구축
python rag/build_index.py \\
  --golden_modules ./golden_modules/ \\
  --hal_sources ./STM32CubeG4/Drivers/ \\
  --datasheet_chunks ./data/chunks/ \\
  --collection stm32g4_knowledge

# 임베딩 모델: BAAI/bge-m3
# 한국어+영어 혼재 환경 최적
# 인덱싱 소요: 약 30~60분"""
add_rect(s, 6.65, 1.82, 6.3, 3.9, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, qdrant_code, 6.8, 1.87, 6.1, 3.8, size=10.5, color=RGBColor(0xAD,0xD8,0xE6))

# 디렉토리 구조
add_text(s, "작업 디렉토리 구조", 0.35, 5.82, 12.6, 0.32, size=13, bold=True, color=TEAL)
dir_text = ("/workspace/stm32_agent/  ├─ data/step1_verification/  train.jsonl · val.jsonl · test.jsonl  "
            "│  └─ step3_integration/  train.jsonl ...  ├─ models/  base/ · finetuned/ · gguf/  "
            "├─ golden_modules/  *.c · *.h  └─ rag/  build_index.py · query.py")
add_rect(s, 0.35, 6.18, 12.6, 0.72, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, dir_text, 0.5, 6.22, 12.3, 0.65, size=10.5, color=RGBColor(0xCE,0xF5,0xB0))

# ══════════════════════════════════════════════════════════════
# B-2: Step 1 QLoRA 학습 하이퍼파라미터
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "B-2  Step 1 모델 — QLoRA Fine-tuning",
             "베이스: Gemma-4-31B-It  |  목적: 검증 리포트 포맷 통일")
footer(s, 30)

add_text(s, "핵심 하이퍼파라미터", 0.35, 1.42, 5.8, 0.35, size=14, bold=True, color=BLUE)
hp_rows = [
    ("LoRA rank (r)",            "32",          "높을수록 품질↑ 메모리↑"),
    ("LoRA alpha",               "32",          "보통 r과 동일하게 설정"),
    ("LoRA dropout",             "0.05",        "과적합 방지"),
    ("Target modules",           "q/k/v/o/gate/up/down", "Gemma-4 전체 레이어"),
    ("Learning rate",            "2e-4",        "cosine 스케줄러"),
    ("Warmup ratio",             "0.05",        "전체 스텝의 5%"),
    ("Batch size",               "2 × GA 8 = 16", "실효 배치 크기"),
    ("Epochs",                   "3",           "검증셋 모니터링 필수"),
    ("Max seq length",           "4096 토큰",   "핀 JSON + 설명 충분"),
    ("Quantization",             "4bit (NF4)",  "학습 중 양자화"),
    ("Precision",                "bfloat16",    "GB10 Blackwell 지원"),
    ("예상 학습 시간 (1,000쌍)", "~4시간",       "DGX Spark 128GB 기준"),
]
add_rect(s, 0.35, 1.82, 5.8, 0.32, fill=NAVY)
add_text(s, "파라미터", 0.45, 1.85, 1.8, 0.25, size=11, bold=True, color=WHITE)
add_text(s, "값",       2.35, 1.85, 1.6, 0.25, size=11, bold=True, color=WHITE)
add_text(s, "설명",     4.05, 1.85, 2.0, 0.25, size=11, bold=True, color=WHITE)
for i, (k, v, desc) in enumerate(hp_rows):
    bg = RGBColor(0xF0,0xF5,0xFF) if i%2==0 else WHITE
    add_rect(s, 0.35, 2.16+i*0.36, 5.8, 0.34, fill=bg)
    add_text(s, k,    0.45, 2.19+i*0.36, 1.8, 0.28, size=11, color=NAVY)
    add_text(s, v,    2.35, 2.19+i*0.36, 1.6, 0.28, size=11, bold=True, color=BLUE)
    add_text(s, desc, 4.05, 2.19+i*0.36, 2.0, 0.28, size=10.5, color=GRAY)

add_text(s, "학습 스크립트 핵심", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=BLUE)
train_code = """from unsloth import FastLanguageModel
from trl import SFTTrainer, SFTConfig

model, tokenizer = FastLanguageModel.from_pretrained(
    model_name = "/workspace/models/base/Gemma-4-31B",
    max_seq_length = 4096,
    load_in_4bit   = True,
)
model = FastLanguageModel.get_peft_model(
    model,
    r=32, lora_alpha=32,
    lora_dropout=0.05,
    target_modules=["q_proj","k_proj","v_proj",
                    "o_proj","gate_proj",
                    "up_proj","down_proj"],
)
trainer = SFTTrainer(
    model=model, tokenizer=tokenizer,
    train_dataset=dataset["train"],
    eval_dataset =dataset["validation"],
    args=SFTConfig(
        per_device_train_batch_size=2,
        gradient_accumulation_steps=8,
        num_train_epochs=3,
        learning_rate=2e-4,
        bf16=True,
        lr_scheduler_type="cosine",
        output_dir="/workspace/models/finetuned/step1",
    )
)
trainer.train()"""
add_rect(s, 6.65, 1.82, 6.3, 4.85, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, train_code, 6.8, 1.87, 6.1, 4.75, size=10, color=RGBColor(0xCE,0xF5,0xB0))

add_rect(s, 0.35, 6.5, 12.6, 0.42, fill=LIGHTBLUE)
add_text(s, "TensorBoard 모니터링:  tensorboard --logdir=/workspace/models/logs  →  http://dgx-spark:6006",
         0.5, 6.54, 12.3, 0.34, size=12, color=NAVY, bold=True)

# ══════════════════════════════════════════════════════════════
# B-3: Step 3 QLoRA + 어댑터 병합 + Ollama 배포
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "B-3  Step 3 모델 — QLoRA Fine-tuning & Ollama 배포",
             "베이스: Gemma-4-31B-It  |  목적: USER CODE 통합 코드 생성")
footer(s, 31)

add_text(s, "Step 1과 다른 설정값", 0.35, 1.42, 5.8, 0.35, size=14, bold=True, color=GREEN)
diff_rows = [
    ("LoRA rank (r)",    "64",              "코드 생성은 rank 높게"),
    ("Learning rate",   "1e-4",            "Step1(2e-4)보다 낮게"),
    ("Max seq length",  "8192 토큰",       "코드가 길어 컨텍스트 2배"),
    ("Batch size",      "1 × GA 16 = 16",  "긴 시퀀스로 배치 줄임"),
    ("Epochs",          "5",               "코드는 더 많은 학습 필요"),
    ("예상 학습 (200쌍)","~6시간",          "DGX Spark 128GB 기준"),
]
add_rect(s, 0.35, 1.82, 5.8, 0.32, fill=NAVY)
add_text(s, "파라미터", 0.45, 1.85, 1.9, 0.25, size=11, bold=True, color=WHITE)
add_text(s, "값",       2.45, 1.85, 1.5, 0.25, size=11, bold=True, color=WHITE)
add_text(s, "Step1 대비", 4.05,1.85, 2.0, 0.25, size=11, bold=True, color=WHITE)
for i, (k, v, desc) in enumerate(diff_rows):
    bg = LIGHTGREEN if i%2==0 else WHITE
    add_rect(s, 0.35, 2.16+i*0.42, 5.8, 0.4, fill=bg)
    add_text(s, k,    0.45, 2.2+i*0.42, 1.9, 0.32, size=11.5, color=NAVY)
    add_text(s, v,    2.45, 2.2+i*0.42, 1.5, 0.32, size=12, bold=True, color=GREEN)
    add_text(s, desc, 4.05, 2.2+i*0.42, 2.0, 0.32, size=11, color=GRAY)

add_text(s, "어댑터 병합 → GGUF 변환 → Ollama 등록", 6.65, 1.42, 6.3, 0.35,
         size=13, bold=True, color=GREEN)
deploy_code = """# 1. 어댑터 + 베이스 모델 병합 후 GGUF 변환
from unsloth import FastLanguageModel
model, tokenizer = FastLanguageModel.from_pretrained(
    "/workspace/models/base/Gemma-4-31B",
    load_in_4bit=True
)
model.load_adapter(
    "/workspace/models/finetuned/step3/adapter"
)
model.save_pretrained_gguf(
    "/workspace/models/gguf/step3",
    tokenizer,
    quantization_method="q8_0"  # Q8 품질 유지
)

# 2. Ollama Modelfile 작성
cat > Modelfile << 'EOF'
FROM /workspace/models/gguf/step3/model-q8_0.gguf
PARAMETER num_ctx 8192
PARAMETER temperature 0.1
SYSTEM "STM32G4 펌웨어 전문가..."
EOF

# 3. Ollama 등록 & 테스트
ollama create stm32-coder-v1 -f Modelfile
ollama run stm32-coder-v1 "BLDC FOC 통합 코드 작성"

# 4. API 서버에서 모델 지정
# POST http://dgx-spark:11434/api/chat
# {"model": "stm32-coder-v1", ...}"""
add_rect(s, 6.65, 1.82, 6.3, 5.08, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, deploy_code, 6.8, 1.87, 6.1, 4.98, size=10, color=RGBColor(0xCE,0xF5,0xB0))

add_rect(s, 0.35, 4.62, 5.8, 0.42, fill=LIGHTGREEN)
add_text(s, "변환 소요:  32B GGUF 약 20~30분  |  등록 후 즉시 API 사용 가능",
         0.5, 4.66, 5.6, 0.34, size=12, color=GREEN, bold=True)

add_rect(s, 0.35, 5.1, 5.8, 0.85, fill=LIGHTORANGE)
add_text(s, "버전 관리 권장", 0.5, 5.14, 5.6, 0.28, size=12, bold=True, color=ORANGE)
add_text(s, "stm32-coder-v1  (기본 RAG 버전)\nstm32-coder-v2  (fine-tuned 50쌍)\nstm32-coder-v3  (fine-tuned 500쌍)",
         0.5, 5.42, 5.6, 0.5, size=11.5, color=BLACK)

# ══════════════════════════════════════════════════════════════
# B-4: 평가 지표
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "B-4  모델 평가 지표 & 합격 기준")
footer(s, 32)

add_text(s, "Step 1 — 핀 검증 평가", 0.35, 1.42, 6.0, 0.35, size=14, bold=True, color=BLUE)
s1_metrics = [
    ("JSON 형식 유효",      "LLM 출력이 파싱 가능한 JSON",      "목표 100%",  "> 99%", GREEN),
    ("오류 Recall",        "실제 오류 핀 중 감지 비율",          "목표 > 95%", "> 90%", GREEN),
    ("False Positive",     "정상 핀을 오류로 판단하는 비율",     "목표 < 5%",  "< 10%", ORANGE),
    ("Exact Match",        "규칙 엔진 정답과 완전 일치",         "목표 > 90%", "> 85%", GREEN),
    ("한국어 설명 품질",    "오류 메시지 자연스러움 (수동 평가)", "목표 > 4/5", "> 3/5", ORANGE),
]
add_rect(s, 0.35, 1.82, 6.0, 0.32, fill=NAVY)
for j, h in enumerate(["지표", "설명", "목표", "최소"]):
    add_text(s, h, 0.45+j*1.48, 1.85, 1.4, 0.25, size=11, bold=True, color=WHITE)
for i, (name, desc, target, minimum, color) in enumerate(s1_metrics):
    bg = RGBColor(0xF0,0xF5,0xFF) if i%2==0 else WHITE
    add_rect(s, 0.35, 2.16+i*0.48, 6.0, 0.46, fill=bg)
    add_text(s, name,    0.45, 2.2+i*0.48, 1.42, 0.38, size=11.5, bold=True, color=NAVY)
    add_text(s, desc,    1.93, 2.2+i*0.48, 1.42, 0.38, size=11, color=GRAY)
    add_text(s, target,  3.41, 2.2+i*0.48, 1.42, 0.38, size=12, bold=True, color=color)
    add_text(s, minimum, 4.89, 2.2+i*0.48, 1.3,  0.38, size=11.5, color=GRAY)

add_text(s, "Step 3 — 코드 통합 평가", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=GREEN)
s3_metrics = [
    ("컴파일 통과",         "arm-none-eabi-gcc 빌드 성공",         "목표 > 95%", "> 90%", GREEN),
    ("HAL 코드 미수정",    "MX_* 함수 원본 유지 여부",             "목표 100%",  "100%",  GREEN),
    ("USER CODE 준수",     "지정 영역에만 코드 삽입",              "목표 100%",  "> 98%", GREEN),
    ("모듈 include 완전",  "필요 헤더 전부 include",               "목표 > 98%", "> 95%", GREEN),
    ("FOC 루프 연결",       "ADC ISR → FOC → TIM 흐름 정확",      "목표 > 90%", "> 85%", ORANGE),
    ("HW 동작 검증",        "실제 모터 회전 확인 (샘플 20%)",      "목표 > 80%", "> 70%", ORANGE),
]
add_rect(s, 6.65, 1.82, 6.3, 0.32, fill=NAVY)
for j, h in enumerate(["지표", "설명", "목표", "최소"]):
    add_text(s, h, 6.75+j*1.55, 1.85, 1.48, 0.25, size=11, bold=True, color=WHITE)
for i, (name, desc, target, minimum, color) in enumerate(s3_metrics):
    bg = LIGHTGREEN if i%2==0 else WHITE
    add_rect(s, 6.65, 2.16+i*0.48, 6.3, 0.46, fill=bg)
    add_text(s, name,    6.75, 2.2+i*0.48, 1.48, 0.38, size=11.5, bold=True, color=GREEN)
    add_text(s, desc,    8.3,  2.2+i*0.48, 1.48, 0.38, size=11, color=GRAY)
    add_text(s, target,  9.85, 2.2+i*0.48, 1.48, 0.38, size=12, bold=True, color=color)
    add_text(s, minimum, 11.4, 2.2+i*0.48, 1.4,  0.38, size=11.5, color=GRAY)

add_rect(s, 0.35, 4.58, 12.6, 0.55, fill=NAVY)
add_text(s, "자동화 평가 파이프라인:  python evaluate.py → compile_test → metrics_report.json → TensorBoard 시각화",
         0.5, 4.64, 12.3, 0.42, size=12.5, bold=True, color=WHITE)
add_rect(s, 0.35, 5.17, 12.6, 0.72, fill=LIGHTTEAL)
add_text(s, "평가 주기 권장:  학습 중 eval_steps=100마다 자동 평가  |  "
         "버전 릴리즈 전 전체 test.jsonl 평가  |  HW 검증은 릴리즈 후 2주 이내 샘플링",
         0.5, 5.22, 12.3, 0.62, size=12, color=TEAL)

# ══════════════════════════════════════════════════════════════
# 어펜딕스 섹션 표지 — C
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
add_rect(s, 0, 0, 13.33, 7.5, fill=NAVY)
add_rect(s, 0, 2.5, 13.33, 2.5, fill=GREEN)
add_text(s, "Appendix C", 1.0, 1.3, 11.0, 0.9,
         size=22, bold=True, color=RGBColor(0xA8,0xD0,0xFF), align=PP_ALIGN.CENTER)
add_text(s, "웹 애플리케이션 개발 프로세스", 1.0, 2.65, 11.0, 0.95,
         size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "DGX Spark 서버 기반  |  HW / FW 개발자 모두 사용", 1.0, 3.7, 11.0, 0.6,
         size=18, color=WHITE, align=PP_ALIGN.CENTER, italic=True)
add_rect(s, 0, 6.9, 13.33, 0.6, fill=RGBColor(0x0F,0x24,0x45))
add_text(s, "FastAPI + React 18 + Tailwind CSS  |  Docker Compose 단일 명령 배포",
         0.5, 6.93, 12.3, 0.4, size=11.5, color=RGBColor(0x99,0xBB,0xFF), align=PP_ALIGN.CENTER)
footer(s, 33)

# ══════════════════════════════════════════════════════════════
# C-1: 기술 스택 & 시스템 아키텍처
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "C-1  기술 스택 & 시스템 아키텍처")
footer(s, 34)

# MVP vs 프로덕션 표
add_text(s, "단계별 기술 스택 선택", 0.35, 1.42, 12.6, 0.35, size=14, bold=True, color=GREEN)
tech_rows = [
    ("프론트엔드",   "Streamlit",             "React 18 + TypeScript"),
    ("UI 라이브러리","Streamlit 기본",         "Tailwind CSS + shadcn/ui"),
    ("백엔드",       "FastAPI",               "FastAPI (동일)"),
    ("실시간 통신",  "폴링 (3초 간격)",        "WebSocket"),
    ("파일 저장",    "로컬 /tmp",             "세션 폴더 + MinIO"),
    ("인증",        "없음",                   "사내 LDAP/AD 연동"),
    ("개발 기간",    "2~3주",                 "6~8주"),
]
add_rect(s, 0.35, 1.82, 12.6, 0.32, fill=NAVY)
for j, h in enumerate(["레이어", "MVP (Streamlit)", "프로덕션 (React)"]):
    add_text(s, h, 0.45+j*4.1, 1.85, 3.9, 0.25, size=12, bold=True, color=WHITE)
for i, (layer, mvp, prod) in enumerate(tech_rows):
    bg = LIGHTGREEN if i%2==0 else WHITE
    add_rect(s, 0.35, 2.16+i*0.38, 12.6, 0.36, fill=bg)
    add_text(s, layer, 0.45, 2.19+i*0.38, 3.9, 0.3, size=12, bold=True, color=NAVY)
    add_text(s, mvp,   4.55, 2.19+i*0.38, 3.9, 0.3, size=12, color=GRAY)
    add_text(s, prod,  8.65, 2.19+i*0.38, 3.9, 0.3, size=12, bold=True, color=GREEN)

# 아키텍처 다이어그램 (텍스트)
add_text(s, "배포 아키텍처 (DGX Spark)", 0.35, 4.95, 12.6, 0.32, size=13, bold=True, color=GREEN)
arch_txt = ("[브라우저]  →  Nginx :80/443  →  React :3000  +  FastAPI :8000\n"
            "FastAPI  →  Ollama :11434 (Gemma-4-31B · Gemma-4-31B)\n"
            "         →  Qdrant :6333 (벡터 DB)  →  CubeMX CLI (subprocess)\n"
            "모두 Docker Compose 단일 명령으로 기동:  docker compose up -d")
add_rect(s, 0.35, 5.3, 12.6, 1.0, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, arch_txt, 0.5, 5.35, 12.3, 0.9, size=12, color=RGBColor(0xCE,0xF5,0xB0))

# ══════════════════════════════════════════════════════════════
# C-2: 백엔드 API + 게이트 로직
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "C-2  백엔드 API 설계 (FastAPI)", "검증 게이트 로직 포함")
footer(s, 35)

add_text(s, "API 엔드포인트 목록", 0.35, 1.42, 5.8, 0.35, size=14, bold=True, color=GREEN)
api_rows = [
    ("POST", "/api/v1/validate",     "핀 검증 (Step 1)", "CSV + 요구사항 → 검증 결과 JSON"),
    ("POST", "/api/v1/generate",     "CubeMX 자동화 (Step 2)", "검증 JSON → HAL 코드 (PASS만 허용)"),
    ("WS",   "/ws/integrate/{id}",   "알고리즘 통합 (Step 3)", "CubeMX 코드 → 완성 펌웨어 스트리밍"),
    ("GET",  "/api/v1/session/{id}", "세션 파일 조회", "생성된 파일 트리 반환"),
    ("GET",  "/api/v1/download/{id}","ZIP 다운로드", "전체 프로젝트 ZIP"),
    ("GET",  "/api/v1/history",      "프로젝트 이력", "이전 생성 이력 목록"),
]
add_rect(s, 0.35, 1.82, 5.8, 0.32, fill=NAVY)
for j, h in enumerate(["Method", "경로", "역할"]):
    add_text(s, h, 0.42+j*1.9, 1.85, 1.85, 0.25, size=11, bold=True, color=WHITE)
for i, (method, path, role, desc) in enumerate(api_rows):
    mc = {"POST": BLUE, "WS": GREEN, "GET": TEAL}[method]
    bg = RGBColor(0xF5,0xF5,0xF5) if i%2==0 else WHITE
    add_rect(s, 0.35, 2.16+i*0.47, 5.8, 0.45, fill=bg)
    add_rect(s, 0.35, 2.16+i*0.47, 0.75, 0.45, fill=mc)
    add_text(s, method, 0.37, 2.2+i*0.47, 0.72, 0.37, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, path,   1.15, 2.2+i*0.47, 2.3, 0.37, size=10.5, bold=True, color=NAVY)
    add_text(s, role,   3.5,  2.2+i*0.47, 2.5, 0.37, size=10.5, color=GRAY)

add_text(s, "검증 게이트 코드 (핵심)", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=GREEN)
gate_code = """# backend/api/step2.py
from fastapi import HTTPException

@app.post("/api/v1/generate")
async def generate_hal_code(
    validated_json: ValidatedPinJSON,
    session_id: str
):
    # ── 검증 게이트 ──────────────────────
    if validated_json.status != "PASS":
        raise HTTPException(
            status_code=403,
            detail={
                "error": "GATE_BLOCKED",
                "message": "핀 검증을 통과해야 합니다.",
                "errors": validated_json.errors
            }
        )
    # errors[] 비어있는지 2중 확인
    if len(validated_json.errors) > 0:
        raise HTTPException(403, "검증 오류가 있습니다.")
    # ──────────────────────────────────────

    # 게이트 통과 시에만 실행
    ioc_path = generate_ioc(validated_json)
    success  = run_cubemx_cli(ioc_path)

    if not success:
        raise HTTPException(500, "CubeMX 실행 실패")

    files = collect_generated_files(session_id)
    return GenerationResult(
        session_id=session_id,
        files=files,
        status="SUCCESS"
    )"""
add_rect(s, 6.65, 1.82, 6.3, 5.55, fill=RGBColor(0x1E,0x1E,0x1E))
add_text(s, gate_code, 6.8, 1.87, 6.1, 5.45, size=10, color=RGBColor(0xCE,0xF5,0xB0))

# ══════════════════════════════════════════════════════════════
# C-3: 프론트엔드 화면 구성 + 개발 단계
# ══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
slide_header(s, "C-3  프론트엔드 화면 구성 & 개발 단계")
footer(s, 36)

add_text(s, "화면 목록 & 주요 기능", 0.35, 1.42, 6.0, 0.35, size=14, bold=True, color=GREEN)
screens = [
    ("/",              "대시보드",   "최근 프로젝트 카드, 새 프로젝트 버튼"),
    ("/step1",         "핀 검증",   "CSV 드래그앤드롭, 요구사항 폼, PASS/FAIL 즉시 표시"),
    ("/step2",         "코드 생성", "자동 진행 바, 생성 파일 트리 미리보기"),
    ("/step3",         "알고리즘 통합","모듈 선택 체크박스, 코드 스트리밍(Monaco Editor)"),
    ("/result",        "결과",      "코드 전체 뷰어, 탭 파일 전환, ZIP 다운로드"),
    ("/history",       "이력",      "이전 프로젝트 목록, 재실행 버튼"),
]
add_rect(s, 0.35, 1.82, 6.0, 0.3, fill=NAVY)
for j, h in enumerate(["경로", "화면명", "핵심 기능"]):
    add_text(s, h, 0.42+j*2.0, 1.85, 1.9, 0.23, size=11, bold=True, color=WHITE)
for i, (path, name, feat) in enumerate(screens):
    bg = LIGHTGREEN if i%2==0 else WHITE
    add_rect(s, 0.35, 2.14+i*0.47, 6.0, 0.45, fill=bg)
    add_text(s, path, 0.42, 2.18+i*0.47, 1.65, 0.37, size=10.5, bold=True, color=GREEN)
    add_text(s, name, 2.12, 2.18+i*0.47, 1.4,  0.37, size=11.5, bold=True, color=NAVY)
    add_text(s, feat, 3.58, 2.18+i*0.47, 2.7,  0.37, size=11, color=GRAY)

# 개발 단계
add_text(s, "개발 단계별 일정", 6.65, 1.42, 6.3, 0.35, size=14, bold=True, color=GREEN)
dev_phases = [
    ("MVP\n2~3주", [
        "FastAPI 3개 엔드포인트 (검증·생성·통합)",
        "Streamlit Step 1~3 화면",
        "Docker Compose 배포",
        "내부 5명 테스트",
    ], TEAL, LIGHTTEAL),
    ("프로덕션\nWeek 1~2", [
        "Next.js 14 + TypeScript 세팅",
        "Tailwind CSS + shadcn/ui 적용",
        "FastAPI WebSocket 추가",
        "라우팅 및 레이아웃",
    ], BLUE, LIGHTBLUE),
    ("프로덕션\nWeek 3~6", [
        "Step 1 드래그앤드롭 + 실시간 결과",
        "Step 2 진행 바 (WebSocket)",
        "Step 3 Monaco Editor 스트리밍",
        "ZIP 다운로드 + 이력 저장",
    ], GREEN, LIGHTGREEN),
    ("프로덕션\nWeek 7~8", [
        "사내 LDAP/AD 인증 연동",
        "Nginx + HTTPS 설정",
        "전체 사용자 테스트",
        "운영 가이드 문서 작성",
    ], ORANGE, LIGHTORANGE),
]
for i, (phase, tasks, color, bg) in enumerate(dev_phases):
    ty = 1.82 + i * 1.33
    add_rect(s, 6.65, ty, 1.05, 1.28, fill=color)
    add_text(s, phase, 6.67, ty+0.25, 1.02, 0.78, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, 7.73, ty, 5.22, 1.28, fill=bg)
    for j, task in enumerate(tasks):
        add_text(s, f"• {task}", 7.88, ty+0.06+j*0.3, 5.0, 0.28, size=11.5, color=BLACK)

# ── 저장 ───────────────────────────────────────────────────
output_path = "/home/younlea/source-code/MotorDriveForge/STM32G4_Agent_Plan.pptx"
prs.save(output_path)
print(f"PPT 생성 완료: {output_path}")